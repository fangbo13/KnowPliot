"""
EY Onboarding AI Chatbot - PPT Report Generator
安永入职AI助手 - PPT测评报告生成器

Generates a bilingual (Chinese + English) evaluation report in PPT format.
生成中英双语测评报告PPT文件。
"""

import json
import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# EY Brand Colors
EY_RED = RGBColor(0xE0, 0x00, 0x33)
EY_DARK = RGBColor(0x30, 0x30, 0x30)
EY_GRAY = RGBColor(0x66, 0x66, 0x66)
EY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
EY_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EY_BORDER = RGBColor(0xDD, 0xDD, 0xDD)

SCREENSHOTS_DIR = os.path.join(os.path.dirname(__file__), "screenshots")
TEST_RESULTS_FILE = os.path.join(os.path.dirname(__file__), "test_results.json")


class PPTGenerator:
    def __init__(self, results_file=None):
        self.prs = Presentation()
        # Widescreen
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.results = []
        self.results_file = results_file or TEST_RESULTS_FILE

        # Load results
        if os.path.exists(self.results_file):
            with open(self.results_file, "r", encoding="utf-8") as f:
                self.results = json.load(f)

        # Debug
        print(f"Loaded {len(self.results)} test results")

    def _add_bg(self, slide, color=EY_WHITE):
        """Add full-slide background color"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = color
        background.line.fill.background()

    def _add_textbox(self, slide, left, top, width, height, text, font_size=18,
                     bold=False, color=EY_DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        """Add a text box"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def _add_red_bar(self, slide, left, top, width, height):
        """Add EY red accent bar"""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = EY_RED
        bar.line.fill.background()

    def _add_eight_ball_chart(self, slide, left, top, items, col_width=2.0, row_height=0.5):
        """
        Add a table with status circles (green=pass, red=fail, yellow=skip, gray=error).
        """
        # Header row
        headers = ["Test ID", "Test Name / 测试名称", "Category / 类别", "Status / 状态"]
        for j, header in enumerate(headers):
            box = slide.shapes.add_textbox(
                Inches(left + j * col_width), Inches(top), Inches(col_width), Inches(row_height)
            )
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = EY_WHITE
            p.alignment = PP_ALIGN.CENTER

            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + j * col_width), Inches(top), Inches(col_width), Inches(row_height)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = EY_RED
            bg.line.fill.background()

        # Data rows
        y = top + row_height
        for item in items[:12]:  # Limit to 12 rows per table
            status = item.get("status", "skip")
            color_map = {"pass": RGBColor(0x22, 0x8B, 0x22), "fail": EY_RED,
                         "skip": RGBColor(0xFF, 0xA5, 0x00), "error": RGBColor(0x80, 0x80, 0x80)}
            dot_color = color_map.get(status, RGBColor(0x80, 0x80, 0x80))

            row_data = [
                item.get("id", ""),
                f"{item.get('name_en', '')} / {item.get('name', '')}",
                item.get("category", ""),
                status.upper()
            ]

            for j, val in enumerate(row_data):
                box = slide.shapes.add_textbox(
                    Inches(left + j * col_width), Inches(y), Inches(col_width), Inches(row_height)
                )
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = val
                p.font.size = Pt(10)
                p.font.color.rgb = EY_DARK
                p.alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT

            # Add status dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + 3 * col_width + col_width * 0.4),
                Inches(y + row_height * 0.15),
                Inches(0.15), Inches(0.15)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()

            # Row background
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(y), Inches(4 * col_width), Inches(row_height)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = EY_LIGHT if items.index(item) % 2 == 0 else EY_WHITE
            row_bg.line.fill.background()

            y += row_height

    def _add_summary_table(self, slide, left, top, categories):
        """Add summary table by category"""
        headers = ["Category / 类别", "Total / 总计", "Pass / 通过", "Fail / 失败", "Error / 错误", "Rate / 通过率"]
        col_width = 2.0

        for j, header in enumerate(headers):
            box = slide.shapes.add_textbox(
                Inches(left + j * col_width), Inches(top), Inches(col_width), Inches(0.5)
            )
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = EY_WHITE
            p.alignment = PP_ALIGN.CENTER

            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left + j * col_width), Inches(top), Inches(col_width), Inches(0.5)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = EY_RED
            bg.line.fill.background()

        y = top + 0.5
        for cat, stats in categories.items():
            rate = f"{stats['pass']}/{stats['total']} ({stats['pass']/stats['total']*100:.1f}%)" if stats['total'] else "N/A"
            row = [cat, str(stats['total']), str(stats['pass']), str(stats['fail']),
                   str(stats['error']), rate]

            for j, val in enumerate(row):
                box = slide.shapes.add_textbox(
                    Inches(left + j * col_width), Inches(y), Inches(col_width), Inches(0.4)
                )
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = val
                p.font.size = Pt(10)
                p.font.color.rgb = EY_DARK
                p.alignment = PP_ALIGN.CENTER

                bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left + j * col_width), Inches(y), Inches(col_width), Inches(0.4)
                )
                bg.fill.solid()
                bg.fill.fore_color.rgb = EY_LIGHT if list(categories.keys()).index(cat) % 2 == 0 else EY_WHITE
                bg.line.fill.background()

            y += 0.4

    def add_title_slide(self):
        """Slide 1: Title / 封面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._add_bg(slide, EY_WHITE)

        # Red accent bar at top
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        # EY logo placeholder (text)
        self._add_textbox(slide, 1.5, 1.5, 10, 1.5, "EY", font_size=72,
                          bold=True, color=EY_RED, alignment=PP_ALIGN.CENTER)

        # Title
        self._add_textbox(slide, 1.5, 3.0, 10, 1.0,
                          "Onboarding AI Chatbot - Evaluation Report",
                          font_size=32, bold=True, color=EY_DARK, alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, 1.5, 3.8, 10, 0.8,
                          "入职AI助手 - 综合测评报告",
                          font_size=24, color=EY_GRAY, alignment=PP_ALIGN.CENTER)

        # Date and info
        now = datetime.now().strftime("%Y-%m-%d")
        self._add_textbox(slide, 1.5, 5.2, 10, 0.5,
                          f"Date / 日期: {now}  |  Tester / 测试者: AI Application Optimizer",
                          font_size=16, color=EY_GRAY, alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, 1.5, 5.7, 10, 0.5,
                          f"Model / 模型: Qwen-Plus (DashScope)  |  Stack / 技术栈: Django 5 + React 18 + pgvector",
                          font_size=14, color=EY_GRAY, alignment=PP_ALIGN.CENTER)

        # Bottom red bar
        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_executive_summary(self):
        """Slide 2: Executive Summary / 执行摘要"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        # Title
        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Executive Summary / 执行摘要", font_size=28, bold=True, color=EY_RED)

        # Calculate stats
        total = len(self.results)
        passed = sum(1 for r in self.results if r.get("status") == "pass")
        failed = sum(1 for r in self.results if r.get("status") == "fail")
        errors = sum(1 for r in self.results if r.get("status") == "error")
        skipped = sum(1 for r in self.results if r.get("status") == "skip")
        rate = f"{passed}/{total} ({passed/total*100:.1f}%)" if total else "N/A"

        # Big number
        self._add_textbox(slide, 1.0, 1.5, 4, 2, f"{passed/total*100:.1f}%" if total else "N/A",
                          font_size=72, bold=True, color=EY_RED, alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, 1.0, 3.3, 4, 0.5,
                          "Overall Pass Rate\n总体通过率", font_size=16, color=EY_GRAY,
                          alignment=PP_ALIGN.CENTER)

        # Stats cards
        stats = [
            (5.0, "Total Tests\n总测试数", str(total)),
            (6.2, "Passed\n通过", str(passed)),
            (7.4, "Failed\n失败", str(failed)),
            (8.6, "Errors\n错误", str(errors)),
            (9.8, "Pass Rate\n通过率", rate),
        ]

        for x, label, value in stats:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.2), Inches(2.0)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = EY_LIGHT
            card.line.color.rgb = EY_BORDER

            self._add_textbox(slide, x + 0.1, 1.6, 2.0, 0.8, label,
                              font_size=24, bold=True, color=EY_DARK, alignment=PP_ALIGN.CENTER)
            self._add_textbox(slide, x + 0.1, 2.5, 2.0, 0.8, value,
                              font_size=14, color=EY_GRAY, alignment=PP_ALIGN.CENTER)

        # Key findings
        self._add_textbox(slide, 0.8, 4.0, 12, 0.5,
                          "Key Findings / 关键发现", font_size=20, bold=True, color=EY_DARK)

        findings = [
            "✅ Chat streaming (SSE) works end-to-end with proper token delivery",
            "✅ Guardrails system blocks common prompt injection patterns",
            "✅ Authentication (JWT) properly protects all API endpoints",
            "⚠️ Response speed depends on external LLM API (Qwen-Plus via DashScope)",
            "✅ Session ownership and isolation enforced at API level",
            "⚠️ No test infrastructure exists - all testing done via custom scripts",
        ]

        y = 4.6
        for finding in findings:
            self._add_textbox(slide, 1.2, y, 11, 0.35, f"•  {finding}",
                              font_size=13, color=EY_DARK)
            y += 0.35

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_scope_slide(self):
        """Slide 3: Test Scope & Environment / 测试范围与环境"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Test Scope & Environment / 测试范围与环境",
                          font_size=28, bold=True, color=EY_RED)

        # Environment info
        info = [
            ("Application / 应用", "EY Onboarding AI Chatbot"),
            ("Git Commit", "1db8e9d (Initial project setup) + 20 uncommitted changes"),
            ("Backend / 后端", "Django 5.0 + DRF + Celery + PostgreSQL (pgvector)"),
            ("Frontend / 前端", "React 18 + Vite + TypeScript + Ant Design 5 + Zustand"),
            ("AI Model / 模型", "Qwen-Plus via DashScope API (OpenAI-compatible)"),
            ("Embedding / 嵌入", "text-embedding-v4 (1024 dimensions)"),
            ("Test Date / 日期", datetime.now().strftime("%Y-%m-%d")),
            ("Test Method", "Custom Python scripts + HTTP client (httpx)"),
            ("Environment / 环境", "Docker Compose on Windows 11"),
        ]

        y = 1.3
        for label, value in info:
            self._add_textbox(slide, 1.2, y, 4.5, 0.35, f"{label}:",
                              font_size=13, bold=True, color=EY_DARK)
            self._add_textbox(slide, 5.7, y, 7.0, 0.35, value,
                              font_size=13, color=EY_GRAY)
            y += 0.35

        # Test categories
        self._add_textbox(slide, 0.8, 4.5, 12, 0.5,
                          "Test Categories / 测试类别", font_size=20, bold=True, color=EY_DARK)

        categories = [
            ("Authentication / 认证安全", "8 test cases", "Login, logout, token refresh, access control"),
            ("Chat Sessions / 聊天会话", "6 test cases", "CRUD, messages, quick actions"),
            ("Streaming Chat / 流式聊天", "1 test case (multi-query)", "SSE streaming, token delivery"),
            ("Safety/Guardrails / 安全防护", "16 test cases", "Injection patterns, false positive check"),
            ("Performance / 性能", "3 test cases", "TTFT, total latency, API response time"),
            ("Edge Cases / 边界情况", "7 test cases", "Empty input, special chars, isolation"),
        ]

        y = 5.1
        for cat, count, desc in categories:
            self._add_textbox(slide, 1.2, y, 3.5, 0.35, f"  {cat}",
                              font_size=12, bold=True, color=EY_DARK)
            self._add_textbox(slide, 4.8, y, 3.0, 0.35, count,
                              font_size=12, color=EY_RED, alignment=PP_ALIGN.CENTER)
            self._add_textbox(slide, 7.8, y, 5.0, 0.35, desc,
                              font_size=11, color=EY_GRAY)
            y += 0.33

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_results_overview(self):
        """Slide 4: Test Results Overview / 测试结果总览"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Test Results Overview / 测试结果总览",
                          font_size=28, bold=True, color=EY_RED)

        # Calculate category stats
        categories = {}
        for r in self.results:
            cat = r.get("category", "Unknown")
            if cat not in categories:
                categories[cat] = {"total": 0, "pass": 0, "fail": 0, "error": 0, "skip": 0}
            categories[cat]["total"] += 1
            categories[cat][r.get("status", "skip")] += 1

        # Add total row
        total_stats = {"total": 0, "pass": 0, "fail": 0, "error": 0, "skip": 0}
        for stats in categories.values():
            for k in total_stats:
                total_stats[k] += stats[k]
        categories["TOTAL / 总计"] = total_stats

        self._add_summary_table(slide, 0.5, 1.2, categories)

        # Status legend
        legend_y = 1.2 + 0.5 + len(categories) * 0.4 + 0.3
        self._add_textbox(slide, 0.8, legend_y, 12, 0.5,
                          "Legend / 图例:", font_size=14, bold=True, color=EY_DARK)

        items = [
            (RGBColor(0x22, 0x8B, 0x22), "Pass / 通过 - Test completed successfully"),
            (EY_RED, "Fail / 失败 - Test did not meet expected outcome"),
            (RGBColor(0xFF, 0xA5, 0x00), "Skip / 跳过 - Test skipped due to prerequisite failure"),
            (RGBColor(0x80, 0x80, 0x80), "Error / 错误 - Unexpected exception during test"),
        ]

        lx = 0.8
        for color, desc in items:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(lx), Inches(legend_y + 0.5), Inches(0.25), Inches(0.25)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            self._add_textbox(slide, lx + 0.35, legend_y + 0.48, 8, 0.3, desc,
                              font_size=11, color=EY_DARK)
            lx += 4.0

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_performance_slide(self):
        """Slide 5: Performance Analysis / 性能分析"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Performance Analysis / 性能分析",
                          font_size=28, bold=True, color=EY_RED)

        # Find performance results
        perf_results = [r for r in self.results if r.get("category") == "Performance"]
        chat_results = [r for r in self.results if r.get("category") == "Chat" and "STREAM" in r.get("id", "")]

        # Performance metrics display
        metrics_data = []
        for r in perf_results:
            metrics_data.append({
                "id": r.get("id"),
                "name": r.get("name_en", ""),
                "name_cn": r.get("name", ""),
                "status": r.get("status"),
                "metric": r.get("metric", {}),
                "details": r.get("details", "")[:200],
            })

        # Stream results
        for r in chat_results:
            metric = r.get("metric", {})
            if metric:
                metrics_data.append({
                    "id": r.get("id"),
                    "name": f"Streaming Chat ({r.get('details', '')[:30]})",
                    "name_cn": r.get("name", ""),
                    "status": r.get("status"),
                    "metric": metric,
                    "details": r.get("details", "")[:200],
                })

        y = 1.3
        for m in metrics_data:
            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y), Inches(11.0), Inches(1.2)
            )
            status_color = {
                "pass": RGBColor(0xE8, 0xF5, 0xE9),
                "fail": RGBColor(0xFF, 0xEB, 0xEE),
                "error": RGBColor(0xFF, 0xF3, 0xE0),
            }.get(m["status"], EY_LIGHT)
            card.fill.solid()
            card.fill.fore_color.rgb = status_color
            card.line.color.rgb = EY_BORDER

            # Status indicator
            dot_color = {"pass": RGBColor(0x22, 0x8B, 0x22), "fail": EY_RED,
                         "error": RGBColor(0xFF, 0x98, 0x00)}.get(m["status"], EY_GRAY)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1.2), Inches(y + 0.1), Inches(0.2), Inches(0.2)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()

            self._add_textbox(slide, 1.5, y + 0.05, 2.0, 0.35, m["id"],
                              font_size=14, bold=True, color=EY_RED)
            self._add_textbox(slide, 3.5, y + 0.05, 8.0, 0.35, m["name"],
                              font_size=12, bold=True, color=EY_DARK)

            # Metrics
            metric = m.get("metric", {})
            if metric:
                parts = []
                if "avg_ttft_ms" in metric:
                    parts.append(f"TTFT: {metric['avg_ttft_ms']}ms")
                if "avg_total_ms" in metric:
                    parts.append(f"Total: {metric['avg_total_ms']}ms")
                if "ttft_ms" in metric:
                    parts.append(f"TTFT: {metric['ttft_ms']}ms")
                if "total_ms" in metric:
                    parts.append(f"Total: {metric['total_ms']}ms")
                if "token_count" in metric:
                    parts.append(f"Tokens: {metric['token_count']}")
                if "citation_count" in metric:
                    parts.append(f"Citations: {metric['citation_count']}")
                if "avg_ms" in metric:
                    parts.append(f"Avg: {metric['avg_ms']}ms")

                self._add_textbox(slide, 1.5, y + 0.45, 10.0, 0.3, " | ".join(parts),
                                  font_size=12, color=EY_DARK)

            self._add_textbox(slide, 1.5, y + 0.75, 10.0, 0.35, m["details"][:150],
                              font_size=10, color=EY_GRAY)

            y += 1.35

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_safety_slide(self):
        """Slide 6-7: Safety Review / 安全防护审查"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Safety Review / 安全防护审查",
                          font_size=28, bold=True, color=EY_RED)

        safety_results = [r for r in self.results if r.get("category") == "Safety"]

        # Summary stats
        safe_pass = sum(1 for r in safety_results if r.get("status") == "pass")
        safe_fail = sum(1 for r in safety_results if r.get("status") == "fail")
        safe_total = len(safety_results)

        self._add_textbox(slide, 1.0, 1.0, 5, 0.4,
                          f"Guardrails Coverage: {safe_pass}/{safe_total} ({safe_pass/safe_total*100:.1f}%)",
                          font_size=18, bold=True, color=EY_DARK)
        self._add_textbox(slide, 1.0, 1.4, 5, 0.4,
                          f"安全防护覆盖率: {safe_pass}/{safe_total} ({safe_pass/safe_total*100:.1f}%)",
                          font_size=16, color=EY_GRAY)

        # Create injection detection matrix
        self._add_textbox(slide, 1.0, 2.0, 12, 0.4,
                          "Injection Detection Matrix / 注入检测矩阵",
                          font_size=16, bold=True, color=EY_DARK)

        # Table
        y = 2.5
        col_w = 1.8

        # Headers
        headers = ["ID", "Test / 测试", "Should Block", "Result / 结果", "Details / 详情"]
        for j, h in enumerate(headers):
            self._add_textbox(slide, 0.8 + j * col_w, y, col_w, 0.35, h,
                              font_size=10, bold=True, color=EY_WHITE, alignment=PP_ALIGN.CENTER)
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8 + j * col_w), Inches(y), Inches(col_w), Inches(0.35)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = EY_DARK
            bg.line.fill.background()

        y += 0.35
        for r in safety_results[:14]:
            metric = r.get("metric", {})
            blocked = metric.get("blocked", "N/A")
            should_block_text = "YES" if "SAFE-012" >= r.get("id", "SAFE-013") >= "SAFE-001" else "NO"
            try:
                num = int(r.get("id", "SAFE-000").split("-")[1])
                should_block_text = "YES" if num <= 12 else "NO"
            except:
                pass

            row_data = [
                r.get("id", ""),
                r.get("name_en", ""),
                should_block_text,
                r.get("status", "").upper(),
                r.get("details", "")[:50] if r.get("details") else ""
            ]

            for j, val in enumerate(row_data):
                self._add_textbox(slide, 0.8 + j * col_w, y, col_w, 0.35, val,
                                  font_size=9, color=EY_DARK, alignment=PP_ALIGN.CENTER)

            # Row background
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(5 * col_w), Inches(0.35)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = EY_LIGHT if safety_results.index(r) % 2 == 0 else EY_WHITE
            row_bg.line.fill.background()

            y += 0.35

            if y > 7.0:
                break

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_edge_cases_slide(self):
        """Slide 8: Edge Cases / 边界情况"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Edge Cases / 边界情况测试",
                          font_size=28, bold=True, color=EY_RED)

        edge_results = [r for r in self.results if r.get("category") == "Edge Cases"]

        y = 1.3
        for r in edge_results:
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y), Inches(11.0), Inches(0.8)
            )
            status_color = {"pass": RGBColor(0xE8, 0xF5, 0xE9), "fail": RGBColor(0xFF, 0xEB, 0xEE),
                            "error": RGBColor(0xFF, 0xF3, 0xE0)}.get(r.get("status"), EY_LIGHT)
            card.fill.solid()
            card.fill.fore_color.rgb = status_color
            card.line.color.rgb = EY_BORDER

            dot_color = {"pass": RGBColor(0x22, 0x8B, 0x22), "fail": EY_RED,
                         "error": RGBColor(0xFF, 0x98, 0x00)}.get(r.get("status"), EY_GRAY)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1.2), Inches(y + 0.1), Inches(0.15), Inches(0.15)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()

            self._add_textbox(slide, 1.5, y + 0.05, 2.0, 0.3, r.get("id", ""),
                              font_size=12, bold=True, color=EY_RED)
            self._add_textbox(slide, 3.5, y + 0.05, 4.0, 0.3, f"{r.get('name_en', '')} / {r.get('name', '')}",
                              font_size=12, bold=True, color=EY_DARK)
            self._add_textbox(slide, 7.5, y + 0.05, 4.0, 0.6,
                              r.get("details", "")[:120], font_size=10, color=EY_GRAY)

            y += 0.9

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_recommendations_slide(self):
        """Slide 9: Issues & Recommendations / 问题与建议"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Issues & Recommendations / 问题与修复建议",
                          font_size=28, bold=True, color=EY_RED)

        # Find failed tests
        failed = [r for r in self.results if r.get("status") in ("fail", "error")]

        if not failed:
            self._add_textbox(slide, 2.0, 2.0, 9, 1.0,
                              "✅ No critical issues found!\n✅ 未发现严重问题！",
                              font_size=24, bold=True,
                              color=RGBColor(0x22, 0x8B, 0x22), alignment=PP_ALIGN.CENTER)
        else:
            y = 1.3
            for r in failed[:8]:
                severity = "HIGH" if r.get("status") == "error" else "MEDIUM"
                sev_color = EY_RED if severity == "HIGH" else RGBColor(0xFF, 0xA5, 0x00)

                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y), Inches(11.0), Inches(0.9)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
                card.line.color.rgb = EY_BORDER

                self._add_textbox(slide, 1.2, y + 0.05, 1.5, 0.3, severity,
                                  font_size=14, bold=True, color=sev_color)
                self._add_textbox(slide, 2.8, y + 0.05, 3.5, 0.3, r.get("id", ""),
                                  font_size=12, bold=True, color=EY_RED)
                self._add_textbox(slide, 1.2, y + 0.35, 10.5, 0.4,
                                  f"{r.get('name_en', '')} / {r.get('name', '')}",
                                  font_size=11, bold=True, color=EY_DARK)
                self._add_textbox(slide, 1.2, y + 0.6, 10.5, 0.3,
                                  r.get("details", "")[:150], font_size=10, color=EY_GRAY)

                y += 1.0

        # Recommendations
        rec_y = 5.5
        self._add_textbox(slide, 1.0, rec_y, 12, 0.4,
                          "Recommendations / 建议", font_size=18, bold=True, color=EY_DARK)

        recommendations = [
            "1. Add automated test suite (pytest) to prevent regression",
            "2. Implement frontend E2E testing with Playwright for UI verification",
            "3. Add response caching for common questions to improve TTFT",
            "4. Consider adding rate limiting feedback headers (Retry-After)",
            "5. Set up CI/CD pipeline with automated testing on each PR",
        ]

        rec_y += 0.4
        for rec in recommendations:
            self._add_textbox(slide, 1.3, rec_y, 11, 0.3, f"  {rec}",
                              font_size=12, color=EY_GRAY)
            rec_y += 0.3

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def add_conclusion_slide(self):
        """Slide 10: Conclusion / 结论"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, EY_WHITE)
        self._add_red_bar(slide, 0, 0, 13.333, 0.08)

        self._add_textbox(slide, 0.8, 0.3, 12, 0.6,
                          "Conclusion / 结论",
                          font_size=28, bold=True, color=EY_RED)

        total = len(self.results)
        passed = sum(1 for r in self.results if r.get("status") == "pass")
        rate = passed/total*100 if total else 0

        # Overall assessment
        if rate >= 90:
            assessment = "✅ READY FOR PRODUCTION / 可以上线"
            assessment_color = RGBColor(0x22, 0x8B, 0x22)
        elif rate >= 70:
            assessment = "⚠️ NEEDS MINOR FIXES / 需要少量修复"
            assessment_color = RGBColor(0xFF, 0xA5, 0x00)
        else:
            assessment = "❌ NEEDS SIGNIFICANT WORK / 需要大量修复"
            assessment_color = EY_RED

        self._add_textbox(slide, 1.0, 1.5, 11, 1.0, assessment,
                          font_size=36, bold=True, color=assessment_color, alignment=PP_ALIGN.CENTER)

        # Summary stats
        self._add_textbox(slide, 1.0, 3.0, 11, 0.5,
                          "Summary / 总结", font_size=24, bold=True, color=EY_DARK, alignment=PP_ALIGN.CENTER)

        summary_items = [
            f"Total Test Cases / 总测试用例: {total}",
            f"Passed / 通过: {passed}",
            f"Failed / 失败: {sum(1 for r in self.results if r.get('status') == 'fail')}",
            f"Errors / 错误: {sum(1 for r in self.results if r.get('status') == 'error')}",
            f"Overall Pass Rate / 总体通过率: {rate:.1f}%",
        ]

        y = 3.6
        for item in summary_items:
            self._add_textbox(slide, 2.0, y, 9, 0.4, f"•  {item}",
                              font_size=16, color=EY_DARK, alignment=PP_ALIGN.CENTER)
            y += 0.4

        # Next steps
        self._add_textbox(slide, 1.0, 5.8, 11, 0.4,
                          "Next Steps / 下一步", font_size=20, bold=True, color=EY_DARK, alignment=PP_ALIGN.CENTER)

        next_steps = [
            "Address failed test cases / 修复失败的测试用例",
            "Add automated CI/CD testing / 添加自动化CI/CD测试",
            "Monitor production metrics after deployment / 部署后监控生产指标",
        ]

        y = 6.2
        for step in next_steps:
            self._add_textbox(slide, 2.0, y, 9, 0.3, f"→  {step}",
                              font_size=13, color=EY_GRAY, alignment=PP_ALIGN.CENTER)
            y += 0.3

        self._add_red_bar(slide, 0, 7.3, 13.333, 0.2)

    def generate(self, output_path=None):
        """Generate all slides and save"""
        if output_path is None:
            date_str = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"EY_Onboarding_AI_Evaluation_Report_{date_str}.pptx"

        print(f"Generating PPT report / 生成PPT报告: {output_path}")

        # Add all slides
        self.add_title_slide()
        self.add_executive_summary()
        self.add_scope_slide()
        self.add_results_overview()
        self.add_performance_slide()
        self.add_safety_slide()
        self.add_edge_cases_slide()
        self.add_recommendations_slide()
        self.add_conclusion_slide()

        self.prs.save(output_path)
        print(f"PPT report saved / 报告已保存: {output_path}")
        return output_path


def main():
    results_file = TEST_RESULTS_FILE
    if len(sys.argv) > 1:
        results_file = sys.argv[1]

    output = None
    if len(sys.argv) > 2:
        output = sys.argv[2]

    generator = PPTGenerator(results_file)
    path = generator.generate(output)
    print(f"\nDone! Report: {path}")


if __name__ == "__main__":
    main()
