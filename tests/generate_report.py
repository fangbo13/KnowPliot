"""
EY Onboarding AI - UI 测评报告生成器
使用 python-pptx 生成带截图的 PPT 报告
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCREENSHOTS_DIR = Path(__file__).parent / "screenshots"
RESULTS_FILE = SCREENSHOTS_DIR / "test_results.json"
OUTPUT_FILE = Path(__file__).parent.parent / "EY_Onboarding_AI_UI_Evaluation_Report.pptx"

# EY 品牌色
EY_RED = RGBColor(0xE0, 0x00, 0x33)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

# 7 维度评分
SCORES = {
    "视觉与设计": 7,
    "交互与可用性": 8,
    "功能完整性": 7,
    "性能": 6,
    "跨端兼容性": 7,
    "边界与异常": 6,
    "安全性": 5,
}

# 额外发现（通过代码审查发现，自动化测试未覆盖）
EXTRA_FINDINGS = [
    {"category": "视觉设计", "severity": "medium", "title": "样式全部为内联样式",
     "description": "所有组件使用内联 style={{}} 对象，缺乏 CSS 架构，维护困难", "is_positive": False},
    {"category": "视觉设计", "severity": "suggestion", "title": "不支持暗色模式",
     "description": "应用没有暗色模式支持，对夜间使用不友好", "is_positive": False},
    {"category": "交互可用性", "severity": "suggestion", "title": "缺少停止生成按钮",
     "description": "AI 响应流式输出时，用户无法中途停止生成", "is_positive": False},
    {"category": "交互可用性", "severity": "suggestion", "title": "缺少消息复制按钮",
     "description": "消息气泡没有一键复制功能，用户需要手动选择文本", "is_positive": False},
    {"category": "交互可用性", "severity": "suggestion", "title": "缺少消息重新生成按钮",
     "description": "无法对 AI 响应进行重新生成", "is_positive": False},
    {"category": "功能测试", "severity": "suggestion", "title": "国际化未实际使用",
     "description": "i18n 翻译文件已加载但未在组件中调用 useTranslation()，双语切换可能无效", "is_positive": False},
    {"category": "功能测试", "severity": "positive", "title": "流式打字光标动画",
     "description": "流式响应过程中显示闪烁光标动画，提供实时反馈", "is_positive": True},
    {"category": "功能测试", "severity": "positive", "title": "引用系统完善",
     "description": "AI 响应包含来源引用卡片，显示文档名、页码和相关度分数", "is_positive": True},
    {"category": "功能测试", "severity": "positive", "title": "知识库管理功能完整",
     "description": "管理员可上传、重新索引、删除文档，支持文件类型筛选和状态管理", "is_positive": True},
    {"category": "性能", "severity": "suggestion", "title": "缺少代码分割",
     "description": "未使用 React.lazy/Suspense 进行路由级代码分割，首屏加载较重", "is_positive": False},
    {"category": "性能", "severity": "suggestion", "title": "依赖包体积较大",
     "description": "Ant Design 5 全量引入，未做 tree-shaking 优化", "is_positive": False},
    {"category": "性能", "severity": "positive", "title": "SSE 流式响应",
     "description": "使用 Server-Sent Events 实现流式输出，用户无需等待完整响应", "is_positive": True},
    {"category": "跨端兼容", "severity": "suggestion", "title": "缺少 PWA 支持",
     "description": "未配置 Service Worker，不支持离线访问和安装为应用", "is_positive": False},
    {"category": "边界测试", "severity": "suggestion", "title": "缺少网络断开处理",
     "description": "SSE 连接中断后用户界面无明确提示", "is_positive": False},
    {"category": "边界测试", "severity": "suggestion", "title": "缺少 404/500 错误页",
     "description": "未定义全局错误边界页面，API 错误时用户体验不明确", "is_positive": False},
    {"category": "安全性", "severity": "positive", "title": "Guardrails 注入防护",
     "description": "后端有 prompt injection 检测（guardrails.py），防御恶意输入", "is_positive": True},
    {"category": "安全性", "severity": "suggestion", "title": "ReactMarkdown 未限制组件",
     "description": "react-markdown 未配置 allowedComponents，可能渲染不安全的 HTML", "is_positive": False},
]

# 加载自动化测试结果
with open(RESULTS_FILE, "r", encoding="utf-8") as f:
    test_results = json.load(f)

all_findings = test_results["findings"] + EXTRA_FINDINGS


def add_slide(prs, title_text, subtitle_text=""):
    """添加幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = EY_RED
    p.alignment = PP_ALIGN.LEFT

    # 红色分隔线
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_RED
    shape.line.fill.background()

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_GRAY

    return slide


def add_text(slide, text, left, top, width, height, font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """添加文本"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=13, bullet_char="•"):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    return txBox


def add_image(slide, image_path, left, top, width, height):
    """添加图片"""
    path = Path(image_path)
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        add_text(slide, f"[截图缺失: {path.name}]", left, top, width, height, 12, color=LIGHT_GRAY)


def add_score_bar(slide, label, score, left, top, width):
    """添加评分条"""
    # 标签
    add_text(slide, label, left, top, 2.5, 0.3, 13, bold=True)

    # 背景条
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(left + 2.5), Inches(top + 0.03), Inches(width - 2.5), Inches(0.25))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    bg_shape.line.fill.background()

    # 填充分数
    fill_width = (width - 2.5) * score / 10
    if fill_width > 0:
        fill_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Inches(left + 2.5), Inches(top + 0.03), Inches(fill_width), Inches(0.25))
        fill_shape.fill.solid()
        # 根据分数设置颜色
        if score >= 8:
            fill_shape.fill.fore_color.rgb = RGBColor(0x52, 0xC4, 0x1A)  # Green
        elif score >= 6:
            fill_shape.fill.fore_color.rgb = RGBColor(0xFA, 0xAD, 0x14)  # Orange
        else:
            fill_shape.fill.fore_color.rgb = EY_RED
        fill_shape.line.fill.background()

    # 分数文字
    add_text(slide, f"{score}/10", left + 2.6 + fill_width, top - 0.02, 0.8, 0.3, 12, bold=True, color=DARK_GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ==========================================
    # Slide 1: 封面
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # EY 红色装饰条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_RED
    shape.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "EY Onboarding AI"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = EY_RED

    p2 = tf.add_paragraph()
    p2.text = "Chatbot UI 综合测评报告"
    p2.font.size = Pt(32)
    p2.font.color.rgb = WHITE

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "视觉设计 · 交互体验 · 功能测试 · 性能分析 · 跨端兼容 · 安全性 · 边界异常"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    p4 = tf2.add_paragraph()
    p4.text = "\n测评日期: 2026-06-18  |  版本: Version_1  |  工具: Playwright + python-pptx"
    p4.font.size = Pt(13)
    p4.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # ==========================================
    # Slide 2: 目录
    # ==========================================
    slide = add_slide(prs, "目录", "报告结构概览")
    toc_items = [
        "1. 测评概述 — 测试范围、方法与环境",
        "2. 七维度评分总览 — 雷达图与评分分析",
        "3. 视觉与设计 — 品牌一致性、色彩、排版",
        "4. 交互与可用性 — 操作流程、反馈、快捷键",
        "5. 功能测试 — 登录、聊天、引用、历史记录",
        "6. 性能分析 — 加载时间、响应时间",
        "7. 跨端兼容性 — 桌面/平板/手机响应式",
        "8. 边界与异常 — 空输入、超长输入、错误处理",
        "9. 安全性 — 认证、授权、XSS防护",
        "10. 问题汇总与改进建议",
        "11. 优点总结",
    ]
    add_bullet_list(slide, toc_items, 1, 1.5, 10, 5, 16)

    # ==========================================
    # Slide 3: 测评概述
    # ==========================================
    slide = add_slide(prs, "测评概述", "测试范围、方法与运行环境")

    add_text(slide, "测试范围", 0.7, 1.5, 2, 0.4, 16, bold=True, color=EY_RED)
    scope_items = [
        "覆盖 5 个页面: 登录、聊天、历史记录、Profile、知识库管理",
        "测试 6 个快捷操作按钮",
        "SSE 流式响应过程观察",
        "响应式布局: 1920px / 1280px / 768px / 375px",
        "边界测试: 空输入、超长输入 (4000 字符)",
        "安全性: 路由保护、Token 存储、硬编码凭据",
    ]
    add_bullet_list(slide, scope_items, 0.7, 1.9, 5.5, 3, 12)

    add_text(slide, "运行环境", 7, 1.5, 2, 0.4, 16, bold=True, color=EY_RED)
    env_items = [
        "前端: React 18 + Ant Design 5 + Vite",
        "后端: Django 5 + DRF + Celery",
        "数据库: PostgreSQL (pgvector)",
        "LLM: Qwen via DashScope API",
        "测试工具: Playwright 1.60 + Edge (headless)",
        "截图: 16 张高清截图",
    ]
    add_bullet_list(slide, env_items, 7, 1.9, 5.5, 3, 12)

    add_text(slide, "性能数据", 0.7, 5, 2, 0.4, 16, bold=True, color=EY_RED)
    perf_data = test_results.get("performance", {})
    perf_items = [
        f"登录页加载: {perf_data.get('login_page_load', 'N/A')}s",
        f"登录到聊天页: {perf_data.get('login_to_chat', 'N/A')}s",
    ]
    add_bullet_list(slide, perf_items, 0.7, 5.4, 5.5, 1.5, 12)

    # ==========================================
    # Slide 4: 七维度评分总览
    # ==========================================
    slide = add_slide(prs, "七维度评分总览", "满分 10 分")

    y_offset = 1.5
    for label, score in SCORES.items():
        add_score_bar(slide, label, score, 0.7, y_offset, 8)
        y_offset += 0.5

    # 平均分
    avg = sum(SCORES.values()) / len(SCORES)
    add_text(slide, f"综合评分: {avg:.1f}/10", 0.7, y_offset + 0.3, 5, 0.5, 20, bold=True, color=EY_RED)

    # 评分说明
    add_text(slide, "评分说明", 7.5, 1.5, 5, 0.4, 14, bold=True, color=EY_RED)
    score_explanation = [
        "8-10: 优秀，功能完善，体验流畅",
        "6-7:  良好，基本功能正常，有改进空间",
        "4-5:  一般，存在明显短板",
        "1-3:  较差，需要重点优化",
    ]
    add_bullet_list(slide, score_explanation, 7.5, 1.9, 5, 2, 12)

    # 截图展示
    add_image(slide, SCREENSHOTS_DIR / "03_welcome_page.png", 7.5, 3.8, 5.5, 3.2)

    # ==========================================
    # Slide 5: 视觉与设计
    # ==========================================
    slide = add_slide(prs, "视觉与设计", "评分: 7/10 — 品牌一致，但样式管理待优化")

    # 优点
    add_text(slide, "✓ 优点", 0.7, 1.5, 1.5, 0.4, 16, bold=True, color=RGBColor(0x52, 0xC4, 0x1A))
    design_pros = [
        "EY 品牌红色 (#E00033) 贯穿登录页标题、用户消息气泡、发送按钮",
        "使用 Ant Design 5 组件库，视觉风格统一",
        "圆角统一为 6px (borderRadius token)，视觉和谐",
        "用户消息 (红色右对齐) 与 AI 消息 (白色左对齐) 视觉区分清晰",
        "Markdown 富文本渲染 (react-markdown)，支持粗体、列表、链接等格式",
        "自定义滚动条样式，细节打磨到位",
    ]
    add_bullet_list(slide, design_pros, 0.7, 1.9, 6, 3, 11)

    # 改进建议
    add_text(slide, "⚠ 改进建议", 0.7, 5.2, 2, 0.4, 16, bold=True, color=RGBColor(0xFA, 0xAD, 0x14))
    design_cons = [
        "全部使用内联 style={{}} 对象，缺乏 CSS 架构 (无 CSS Modules/Styled Components)",
        "不支持暗色模式，对暗色主题用户不友好",
        "缺少响应式字体大小调整 (font-size 固定)",
    ]
    add_bullet_list(slide, design_cons, 0.7, 5.6, 6, 2, 11)

    # 截图
    add_image(slide, SCREENSHOTS_DIR / "14_responsive_1920.png", 7.2, 1.5, 5.8, 5.5)

    # ==========================================
    # Slide 6: 交互与可用性
    # ==========================================
    slide = add_slide(prs, "交互与可用性", "评分: 8/10 — 操作流畅，反馈及时")

    add_text(slide, "✓ 优点", 0.7, 1.5, 1.5, 0.4, 16, bold=True, color=RGBColor(0x52, 0xC4, 0x1A))
    ux_pros = [
        "登录失败显示明确的错误提示 (红色边框 + 文字)",
        "登录成功后自动跳转到 /chat 页",
        "6 个快捷操作按钮，覆盖 IT Setup / Reimbursement / Annual Leave / Training / Office / Buddy",
        "新消息到达时自动滚动到底部 (scrollIntoView)",
        "SSE 流式响应 + 闪烁光标动画，实时反馈",
        "输入框字符计数 (showCount) + maxLength=4000",
        "空输入时发送按钮自动禁用",
        "流式期间输入框禁用，防止重复发送",
    ]
    add_bullet_list(slide, ux_pros, 0.7, 1.9, 6, 3.5, 11)

    add_text(slide, "⚠ 改进建议", 0.7, 5.6, 2, 0.4, 16, bold=True, color=RGBColor(0xFA, 0xAD, 0x14))
    ux_cons = [
        "缺少「停止生成」按钮 — 用户无法中断 AI 响应",
        "缺少消息「复制」按钮 — 需要手动选择文本",
        "缺少消息「重新生成」功能",
    ]
    add_bullet_list(slide, ux_cons, 0.7, 6, 6, 1.5, 11)

    add_image(slide, SCREENSHOTS_DIR / "13_input_area.png", 7.2, 1.5, 5.8, 5.5)

    # ==========================================
    # Slide 7: 功能测试 — 登录与聊天
    # ==========================================
    slide = add_slide(prs, "功能测试 — 登录与聊天", "评分: 7/10")

    add_text(slide, "✓ 功能正常", 0.7, 1.5, 2, 0.4, 16, bold=True, color=RGBColor(0x52, 0xC4, 0x1A))
    func_pros = [
        "JWT 认证登录 — Token 获取 + 用户信息加载正常",
        "SSE 流式响应 — 服务器逐 token 推送，前端实时渲染",
        "来源引用系统 — 每条 AI 响应附带文档来源卡片 (标题/页码/相关度)",
        "会话管理 — 创建/切换/历史会话正常 (25 个历史会话)",
        "侧边栏导航 — Chat / History / Knowledge Base / Profile",
        "知识库管理 (Admin) — 文档列表、上传、重新索引、删除功能完整",
        "路由保护 — 未认证用户自动重定向到登录页",
    ]
    add_bullet_list(slide, func_pros, 0.7, 1.9, 6, 3.5, 11)

    add_text(slide, "⚠ 待验证/改进", 0.7, 5.6, 2, 0.4, 16, bold=True, color=RGBColor(0xFA, 0xAD, 0x14))
    func_cons = [
        "国际化 (i18n) — 翻译文件已加载但未在组件中调用 useTranslation()，中英切换可能无效",
        "消息反馈 (Feedback) — 后端有 /feedback/ 接口，前端未展示反馈按钮",
        "会话删除 — 后端有 DELETE 接口，前端未展示删除按钮",
    ]
    add_bullet_list(slide, func_cons, 0.7, 6, 6, 1.5, 11)

    # 截图: 登录页 + AI响应
    add_image(slide, SCREENSHOTS_DIR / "01_login_page.png", 7.2, 1.5, 2.8, 2.5)
    add_image(slide, SCREENSHOTS_DIR / "07_action_response.png", 10.2, 1.5, 2.8, 2.5)
    add_image(slide, SCREENSHOTS_DIR / "16_responsive_768_chat.png", 7.2, 4.2, 5.8, 3)

    # ==========================================
    # Slide 8: 功能测试 — 各页面
    # ==========================================
    slide = add_slide(prs, "功能测试 — 各页面", "5 个主要页面功能正常")

    add_image(slide, SCREENSHOTS_DIR / "09_history_page.png", 0.7, 1.5, 3.8, 3)
    add_image(slide, SCREENSHOTS_DIR / "10_profile_page.png", 4.8, 1.5, 3.8, 3)
    add_image(slide, SCREENSHOTS_DIR / "11_knowledge_page.png", 8.9, 1.5, 3.8, 3)

    add_text(slide, "历史记录页 (History)", 0.7, 4.6, 3.8, 0.3, 12, bold=True)
    add_text(slide, "• 显示 25 个历史会话\n• 标题 + 日期\n• 点击恢复会话", 0.7, 4.9, 3.8, 1, 10)

    add_text(slide, "用户设置页 (Profile)", 4.8, 4.6, 3.8, 0.3, 12, bold=True)
    add_text(slide, "• 邮箱显示 (只读)\n• 语言偏好选择 (English/Chinese)\n• 表单提交反馈", 4.8, 4.9, 3.8, 1, 10)

    add_text(slide, "知识库管理页 (Admin)", 8.9, 4.6, 3.8, 0.3, 12, bold=True)
    add_text(slide, "• 文档表格 (标题/分类/状态/创建时间)\n• 上传/重新索引/删除\n• 分页功能", 8.9, 4.9, 3.8, 1, 10)

    # ==========================================
    # Slide 9: 性能分析
    # ==========================================
    slide = add_slide(prs, "性能分析", "评分: 6/10 — 基本流畅，有优化空间")

    perf_data = test_results.get("performance", {})

    # 性能数据表格
    add_text(slide, "性能测量数据", 0.7, 1.5, 3, 0.4, 16, bold=True, color=EY_RED)

    perf_table = [
        ["指标", "测量值", "评级"],
        ["登录页加载 (LCP)", f"{perf_data.get('login_page_load', 'N/A')}s", "良好" if perf_data.get('login_page_load', 99) < 3 else "待优化"],
        ["登录到聊天页", f"{perf_data.get('login_to_chat', 'N/A')}s", "良好" if perf_data.get('login_to_chat', 99) < 5 else "待优化"],
    ]

    table_shape = slide.shapes.add_table(len(perf_table), 3, Inches(0.7), Inches(2), Inches(5), Inches(1.2))
    table = table_shape.table

    for row_idx, row_data in enumerate(perf_table):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = EY_RED
            else:
                p.font.color.rgb = DARK_GRAY

    perf_issues = [
        "首屏加载 2.4s — 可接受但未达理想 (<1.5s)",
        "登录跳转 3.2s — JWT 两次请求 + 页面渲染",
        "Ant Design 5 全量引入，未做 tree-shaking",
        "未使用 React.lazy/Suspense 路由级代码分割",
        "未配置资源预加载 (preload/prefetch)",
    ]
    add_text(slide, "性能分析", 0.7, 3.5, 2, 0.4, 16, bold=True, color=EY_RED)
    add_bullet_list(slide, perf_issues, 0.7, 3.9, 5.5, 2.5, 11)

    add_text(slide, "✓ 性能优点", 0.7, 5.8, 2, 0.4, 14, bold=True, color=RGBColor(0x52, 0xC4, 0x1A))
    perf_pros = [
        "SSE 流式响应 — 首 token 快速到达，用户无需等待完整响应",
        "Vite HMR 热更新 — 开发体验流畅",
    ]
    add_bullet_list(slide, perf_pros, 0.7, 6.1, 5.5, 1, 11)

    add_image(slide, SCREENSHOTS_DIR / "15_responsive_1280.png", 7, 1.5, 6, 5.5)

    # ==========================================
    # Slide 10: 跨端兼容性
    # ==========================================
    slide = add_slide(prs, "跨端兼容性", "评分: 7/10 — 响应式良好，缺少 PWA")

    add_image(slide, SCREENSHOTS_DIR / "04_welcome_tablet.png", 0.7, 1.5, 2.8, 3)
    add_image(slide, SCREENSHOTS_DIR / "05_welcome_mobile.png", 3.7, 1.5, 2.8, 3)
    add_image(slide, SCREENSHOTS_DIR / "16_responsive_768_chat.png", 6.7, 1.5, 2.8, 3)
    add_image(slide, SCREENSHOTS_DIR / "14_responsive_1920.png", 9.7, 1.5, 3.3, 3)

    add_text(slide, "平板 768px", 0.7, 4.6, 2.8, 0.3, 11, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, "手机 375px", 3.7, 4.6, 2.8, 0.3, 11, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, "平板聊天 768px", 6.7, 4.6, 2.8, 0.3, 11, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, "桌面 1920px", 9.7, 4.6, 3.3, 0.3, 11, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    responsive_items = [
        "✓ 侧边栏 breakpoint='lg' — 平板下自动折叠为图标模式",
        "✓ 欢迎页 Grid: xs=24 / sm=12 / md=8 — 响应式列数",
        "✓ 快捷按钮在移动端自动变为全宽单列",
        "✓ 聊天消息区域自适应不同宽度",
        "⚠ 缺少 PWA 支持 (无 Service Worker)",
        "⚠ 未做移动端触摸优化 (无 tap-highlight 处理)",
    ]
    add_bullet_list(slide, responsive_items, 0.7, 5.1, 12, 2.5, 12)

    # ==========================================
    # Slide 11: 边界与异常
    # ==========================================
    slide = add_slide(prs, "边界与异常处理", "评分: 6/10 — 基本防护到位，细节待完善")

    add_image(slide, SCREENSHOTS_DIR / "02_login_error.png", 0.7, 1.5, 4, 3)

    add_text(slide, "✓ 已实现的防护", 5.2, 1.5, 3, 0.4, 14, bold=True, color=RGBColor(0x52, 0xC4, 0x1A))
    edge_pros = [
        "空输入时发送按钮禁用",
        "maxLength=4000 字符限制",
        "字数统计显示 (showCount)",
        "登录错误提示清晰",
    ]
    add_bullet_list(slide, edge_pros, 5.2, 1.9, 7, 2, 12)

    add_text(slide, "⚠ 缺失的处理", 5.2, 4, 3, 0.4, 14, bold=True, color=RGBColor(0xFA, 0xAD, 0x14))
    edge_cons = [
        "SSE 连接中断后无用户提示 (仅 console.log)",
        "后端 500 错误时无友好错误页",
        "网络断开时无离线提示",
        "快速连点发送按钮无防抖处理",
        "未定义全局错误边界 (Error Boundary)",
    ]
    add_bullet_list(slide, edge_cons, 5.2, 4.4, 7, 2.5, 12)

    # ==========================================
    # Slide 12: 安全性
    # ==========================================
    slide = add_slide(prs, "安全性", "评分: 5/10 — 基础认证完善，存在安全隐患")

    add_image(slide, SCREENSHOTS_DIR / "01_login_page.png", 0.7, 1.5, 3, 2.5)

    add_text(slide, "✓ 安全措施", 4.2, 1.5, 3, 0.4, 14, bold=True, color=RGBColor(0x52, 0xC4, 0x1A))
    sec_pros = [
        "JWT Token 认证 (access + refresh)",
        "路由保护 — 未认证自动重定向到登录页",
        "CORS 配置 — 限制允许的来源",
        "后端 Guardrails — prompt injection 检测 (guardrails.py)",
        "Token 黑名单机制 — 登出后 token 失效",
    ]
    add_bullet_list(slide, sec_pros, 4.2, 1.9, 8, 2.5, 12)

    add_text(slide, "⚠ 安全隐患", 4.2, 4.5, 3, 0.4, 14, bold=True, color=RGBColor(0xFA, 0xAD, 0x14))
    sec_cons = [
        "登录页硬编码 demo 凭据 (admin@ey.com / admin123) — 生产环境严重风险",
        "JWT Token 存储于 localStorage — XSS 攻击可窃取 token，建议改用 httpOnly cookie",
        "ReactMarkdown 未配置 allowedComponents — 可能渲染不安全的 HTML",
        "密码明文显示在 input[type='password'] — 登录页直接显示预填密码",
    ]
    add_bullet_list(slide, sec_cons, 4.2, 4.9, 8, 2.5, 12)

    # ==========================================
    # Slide 13: 问题汇总
    # ==========================================
    slide = add_slide(prs, "问题汇总", "按严重程度排列")

    issues = [f for f in all_findings if not f["is_positive"]]
    issues.sort(key=lambda x: {"critical": 0, "medium": 1, "suggestion": 2}.get(x["severity"], 3))

    # 严重/中等
    add_text(slide, "🔴 中等风险 (需尽快处理)", 0.7, 1.5, 6, 0.4, 14, bold=True, color=EY_RED)
    y = 1.9
    for issue in issues:
        if issue["severity"] in ["critical", "medium"]:
            sev_icon = "🔴" if issue["severity"] == "critical" else "🟠"
            add_text(slide, f"{sev_icon} [{issue['category']}] {issue['title']}", 0.7, y, 12, 0.3, 12, bold=True)
            add_text(slide, f"   {issue['description']}", 0.7, y + 0.25, 12, 0.3, 10, color=LIGHT_GRAY)
            y += 0.6

    add_text(slide, "🟡 建议 (可后续优化)", 0.7, y + 0.2, 6, 0.4, 14, bold=True, color=RGBColor(0xFA, 0xAD, 0x14))
    y += 0.6
    for issue in issues:
        if issue["severity"] == "suggestion":
            add_text(slide, f"• [{issue['category']}] {issue['title']}", 0.7, y, 12, 0.3, 11)
            y += 0.35

    # ==========================================
    # Slide 14: 改进建议与优先级
    # ==========================================
    slide = add_slide(prs, "改进建议与优先级", "按影响/投入排序")

    recommendations = [
        {"priority": "P0", "title": "移除硬编码 demo 凭据",
         "desc": "登录页不应预填真实账号密码，应使用环境变量配置", "effort": "低"},
        {"priority": "P0", "title": "Token 改用 httpOnly cookie",
         "desc": "localStorage 存储 JWT 有 XSS 风险，改用 httpOnly cookie 更安全", "effort": "中"},
        {"priority": "P1", "title": "启用国际化 (useTranslation)",
         "desc": "组件中调用 useTranslation() 以真正支持中英双语", "effort": "中"},
        {"priority": "P1", "title": "添加代码分割",
         "desc": "使用 React.lazy + Suspense 对路由进行代码分割", "effort": "低"},
        {"priority": "P2", "title": "添加停止生成/复制/重新生成按钮",
         "desc": "提升聊天交互体验", "effort": "中"},
        {"priority": "P2", "title": "完善错误边界和离线处理",
         "desc": "添加 Error Boundary、网络断开提示、SSE 中断恢复", "effort": "中"},
        {"priority": "P3", "title": "CSS 架构优化",
         "desc": "将内联样式迁移到 CSS Modules 或 Styled Components", "effort": "高"},
        {"priority": "P3", "title": "添加暗色模式支持",
         "desc": "使用 Ant Design 的暗色主题 token", "effort": "中"},
    ]

    for i, rec in enumerate(recommendations):
        row = 1.5 + i * 0.65
        # Priority badge
        color = EY_RED if rec["priority"] == "P0" else RGBColor(0xFA, 0xAD, 0x14) if rec["priority"] == "P1" else RGBColor(0x52, 0xC4, 0x1A)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.7), Inches(row), Inches(0.5), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = rec["priority"]
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text(slide, rec["title"], 1.4, row - 0.02, 4, 0.2, 13, bold=True)
        add_text(slide, f"{rec['desc']}  (投入: {rec['effort']})", 1.4, row + 0.2, 10, 0.3, 10, color=LIGHT_GRAY)

    # ==========================================
    # Slide 15: 优点总结
    # ==========================================
    slide = add_slide(prs, "优点总结", f"共发现 {len([f for f in all_findings if f['is_positive']])} 项优点")

    positives = [f for f in all_findings if f["is_positive"]]
    positives.sort(key=lambda x: x["category"])

    categories = {}
    for p in positives:
        cat = p["category"]
        if cat not in categories:
            categories[cat] = []
        categories[cat].append(p["title"])

    y = 1.5
    for cat, items in categories.items():
        add_text(slide, f"▎{cat}", 0.7, y, 3, 0.4, 14, bold=True, color=EY_RED)
        y += 0.4
        for item in items:
            add_text(slide, f"✓ {item}", 0.7, y, 12, 0.3, 12)
            y += 0.3
        y += 0.15

    # 右侧放截图
    add_image(slide, SCREENSHOTS_DIR / "08_sidebar.png", 8, 1.5, 5, 3)
    add_image(slide, SCREENSHOTS_DIR / "12_message_bubbles.png", 8, 4.7, 5, 2.5)

    # ==========================================
    # Slide 16: 结束页
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_RED
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "感谢阅读"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "EY Onboarding AI — UI 综合测评报告"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\n测评工具: Playwright + python-pptx  |  截图: 16 张  |  发现: 28 项"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p3.alignment = PP_ALIGN.CENTER

    # 保存
    prs.save(str(OUTPUT_FILE))
    print(f"\n[OK] PPT report generated: {OUTPUT_FILE}")
    print(f"  幻灯片数: {len(prs.slides)}")
    print(f"  发现数: {len(all_findings)} (优点: {len([f for f in all_findings if f['is_positive']])}, 问题: {len([f for f in all_findings if not f['is_positive']])})")


if __name__ == "__main__":
    main()
