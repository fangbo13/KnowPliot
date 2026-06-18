"""
EY Onboarding AI - UI Premium Optimization Report (Round 2)
Comprehensive UI/UX overhaul: design tokens, layout fixes, premium polish, responsive design.
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import datetime

SS_DIR = Path(__file__).parent / "screenshots" / "ui-optimized"
OUTPUT_FILE = Path(__file__).parent.parent / "EY_Onboarding_AI_UI_Premium_Optimization_Report.pptx"

EY_YELLOW = RGBColor(0xFF, 0xE5, 0x00)
EY_BLACK = RGBColor(0x26, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MID = RGBColor(0x8C, 0x8C, 0x8C)
DARK_BG = RGBColor(0x1F, 0x1F, 0x1F)
GREEN = RGBColor(0x52, 0xC4, 0x1A)
RED = RGBColor(0xFF, 0x4D, 0x4F)


def add_bar(slide, left, top, width, height=0.05, color=EY_YELLOW):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, size=14, bold=False, color=EY_BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Plus Jakarta Sans'
    p.alignment = align
    return txBox


def add_bullets(slide, items, left, top, width, height, size=11, color=EY_BLACK, spacing=4):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Plus Jakarta Sans'
        p.space_after = Pt(spacing)
    return txBox


def add_img(slide, path, left, top, width, height):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(left), Inches(top), Inches(width), Inches(height))
        return True
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GRAY_LIGHT
        shape.line.color.rgb = GRAY_MID
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f'Image:\n{p.name}'
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY_MID
        p.alignment = PP_ALIGN.CENTER
        return False


def add_card(slide, left, top, width, height, title, value, desc, accent_color=EY_YELLOW):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    add_bar(slide, left, top, width, 0.06, accent_color)
    add_text(slide, value, left + 0.15, top + 0.1, width - 0.3, 0.5, 28, True, accent_color)
    add_text(slide, title, left + 0.15, top + 0.6, width - 0.3, 0.3, 12, True)
    add_text(slide, desc, left + 0.15, top + 0.9, width - 0.3, height - 1.0, 9, False, GRAY_MID)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GRAY_LIGHT
    add_bar(slide, 0, 0, 13.33, 0.08)

    # Logo
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(1.6), Inches(2.5), Inches(1.0))
    logo.fill.solid()
    logo.fill.fore_color.rgb = EY_YELLOW
    logo.line.fill.background()
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = 'EY'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = EY_BLACK
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, 'EY Onboarding AI', 1.5, 2.9, 10.33, 0.6, 36, True, EY_BLACK, PP_ALIGN.CENTER)
    add_text(slide, 'UI Premium Optimization Report', 2, 3.5, 9.33, 0.5, 24, False, EY_BLACK, PP_ALIGN.CENTER)
    add_text(slide, 'Comprehensive UI/UX Overhaul: Design Tokens · Layout Fixes · Premium Polish · Responsive Design',
             2.5, 4.2, 8.33, 0.6, 13, False, GRAY_MID, PP_ALIGN.CENTER)
    add_text(slide, f'Date: {datetime.date.today().strftime("%Y-%m-%d")}  |  Version: 2.0  |  Files Modified: 10',
             4.5, 5.0, 4.33, 0.4, 11, False, GRAY_MID, PP_ALIGN.CENTER)

    # ========== SLIDE 2: Table of Contents ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, 'Report Contents', 0.8, 0.4, 6, 0.6, 28, True)

    sections = [
        ('01', 'Optimization Overview', 'Scope, objectives, and design direction'),
        ('02', 'Design Token Overhaul', 'Typography, shadows, noise texture, gradients'),
        ('03', 'Critical Bug Fixes', '12 overflow, responsive, and validation issues resolved'),
        ('04', 'Login Page', 'Responsive brand panel, premium styling, page-load animation'),
        ('05', 'Welcome Screen', 'Premium card buttons, logo upgrade, hover micro-interactions'),
        ('06', 'Chat Page', 'Input area redesign, streaming fixes, thinking indicator'),
        ('07', 'Message Bubbles', 'Overflow fixes, cursor polish, hover elevation, entrance animations'),
        ('08', 'Profile & Dark Mode', 'Dark mode toggle, refined tokens, responsive layout'),
        ('09', 'Knowledge Base & History', 'Table scroll, upload validation, delete confirmation, empty state'),
        ('10', 'Responsive Design', 'Desktop (1440px), Tablet (768px), Mobile (375px)'),
        ('11', 'Summary & Metrics', 'Before/after comparison and improvements'),
    ]
    for i, (num, title, desc) in enumerate(sections):
        y = 1.3 + i * 0.52
        add_text(slide, num, 1.2, y, 0.6, 0.4, 18, True, EY_YELLOW)
        add_text(slide, title, 1.9, y, 4, 0.3, 14, True)
        add_text(slide, desc, 1.9, y + 0.25, 8, 0.25, 10, False, GRAY_MID)

    # ========== SLIDE 3: Overview ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '01  Optimization Overview', 0.8, 0.4, 10, 0.6, 28, True)

    add_text(slide, 'Why Optimize?', 0.8, 1.3, 5.5, 0.4, 18, True)
    add_bullets(slide, [
        '• Generic system fonts (Arial, Roboto) — lacked brand identity',
        '• Zero responsive design — no media queries, broken on mobile/tablet',
        '• Critical overflow bugs — URLs, citation titles, tables broke containers',
        '• Flat, basic shadows — no visual depth or premium feel',
        '• No micro-interactions — static, no hover states or transitions',
        '• Login brand panel broke on screens < 800px',
        '• Knowledge Base: no table scroll, upload validation, or delete confirm',
    ], 0.8, 1.8, 5.8, 3.5, 11)

    add_text(slide, 'Design Direction', 7, 1.3, 5.5, 0.4, 18, True)
    add_bullets(slide, [
        '• Style: Luxury / Refined Minimalism',
        '• Typography: Plus Jakarta Sans (geometric, modern, professional)',
        '• Colors: EY Yellow (#FFE500) surgical accent + Black (#262626)',
        '• Shadows: Layered premium depth, 4 levels + yellow glow',
        '• Texture: Subtle noise overlay (opacity 0.015)',
        '• Animations: Subtle micro-interactions (fade-in, hover elevation)',
        '• Responsive: Breakpoints at 375px, 576px, 768px, 992px, 1440px',
    ], 7, 1.8, 5.8, 3.5, 11)

    # Scope box
    scope = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.4))
    scope.fill.solid()
    scope.fill.fore_color.rgb = GRAY_LIGHT
    scope.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    add_text(slide, 'Scope: All Pages — 10 Files Modified', 1.0, 5.6, 11, 0.35, 14, True)
    add_bullets(slide, [
        'Login Page  |  Chat Page (Welcome + Messages)  |  Profile Page  |  History Page  |  Knowledge Base Page',
        '12 Critical Bugs Fixed  |  20+ UI Improvements  |  6 New Animations  |  Full Responsive System  |  Dark Mode Toggle',
    ], 1.0, 6.0, 11, 0.8, 11, GRAY_MID)

    # ========== SLIDE 4: Design Tokens ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '02  Design Token Overhaul', 0.8, 0.4, 10, 0.6, 28, True)

    # Typography comparison
    add_text(slide, 'Typography Upgrade', 0.8, 1.2, 5.5, 0.4, 18, True)
    font_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.6))
    font_box.fill.solid()
    font_box.fill.fore_color.rgb = WHITE
    font_box.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    add_text(slide, 'Before: System Fonts', 1.0, 1.85, 5, 0.3, 11, False, GRAY_MID)
    add_text(slide, 'Aa Bb Cc Dd Ee Ff Gg 0123456789', 1.0, 2.15, 5, 0.5, 24, False, GRAY_MID)
    add_text(slide, '-apple-system, BlinkMacSystemFont, Segoe UI, Roboto, Arial', 1.0, 2.65, 5, 0.25, 9, False, GRAY_MID)

    add_text(slide, 'After: Plus Jakarta Sans', 1.0, 3.05, 5, 0.3, 11, True)
    add_text(slide, 'Aa Bb Cc Dd Ee Ff Gg 0123456789', 1.0, 3.35, 5, 0.5, 24, True, EY_BLACK)
    add_text(slide, 'Weights: 300 Light · 400 Regular · 500 Medium · 600 Semibold · 700 Bold', 1.0, 3.85, 5, 0.25, 9, False, GRAY_MID)

    # Shadows
    add_text(slide, 'Layered Premium Shadows', 7, 1.2, 5.5, 0.4, 18, True)
    shadows = [
        ('shadow-sm', '0 1px 3px rgba(0,0,0,0.04), 0 1px 2px rgba(0,0,0,0.06)'),
        ('shadow-md', '0 4px 6px rgba(0,0,0,0.03), 0 2px 4px rgba(0,0,0,0.04)'),
        ('shadow-lg', '0 10px 25px rgba(0,0,0,0.06), 0 6px 12px rgba(0,0,0,0.04)'),
        ('shadow-xl', '0 20px 50px rgba(0,0,0,0.1), 0 10px 20px rgba(0,0,0,0.06)'),
        ('shadow-glow', '0 0 20px rgba(255, 229, 0, 0.15) — yellow accent glow'),
    ]
    for i, (name, val) in enumerate(shadows):
        y = 1.7 + i * 0.48
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(y), Inches(5.8), Inches(0.42))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(min(0xFA + i*2, 0xFF), min(0xFA + i*2, 0xFF), min(0xFA + i*2, 0xFF))
        box.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        add_text(slide, name, 7.15, y + 0.03, 1.6, 0.35, 10, True)
        add_text(slide, val[:65], 8.8, y + 0.03, 3.8, 0.35, 9, False, GRAY_MID)

    # Animations
    add_text(slide, 'New Animations', 0.8, 4.6, 5.5, 0.4, 18, True)
    add_bullets(slide, [
        '• fadeInUp: page/component entrance (0.35s ease-out)',
        '• fadeIn: subtle reveal for streaming indicators',
        '• dotBounce: thinking dots with staggered delays',
        '• shimmer: loading skeleton with gradient sweep',
        '• pulseGlow: glowing pulse for interactive elements',
        '• blink: smooth cursor blink (0.8s ease-in-out)',
    ], 0.8, 5.0, 5.8, 2.0, 11)

    # Background effects
    add_text(slide, 'Background Enhancements', 7, 4.6, 5.5, 0.4, 18, True)
    add_bullets(slide, [
        '• Noise texture overlay via SVG filter (opacity: 0.015)',
        '• Warm gradient: linear-gradient(180deg, rgba(255,229,0,0.02), transparent)',
        '• Firefox scrollbar support: scrollbar-width: thin',
        '• WebKit scrollbar: 6px custom styled with gray thumb',
    ], 7, 5.0, 5.8, 2.0, 11)

    # ========== SLIDE 5: Login Page ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '03  Login Page — Premium Redesign', 0.8, 0.4, 10, 0.6, 28, True)

    add_img(slide, SS_DIR / '01_login_page_1440.png', 0.5, 1.1, 6.5, 4.0)

    add_text(slide, 'Improvements', 7.3, 1.1, 5.5, 0.4, 18, True)
    add_bullets(slide, [
        '✓ Plus Jakarta Sans typography throughout',
        '✓ Refined layered shadows for visual depth',
        '✓ Subtle noise texture background overlay',
        '✓ Warm gradient accent (yellow, 0.02 opacity)',
        '✓ Left-edge yellow gradient accent stripe',
        '✓ Login button with yellow glow on hover',
        '✓ Fade-in animation on load (fadeInUp 0.4s)',
        '✓ Responsive: brand panel hides below 800px',
        '✓ Brand panel 3-stop diagonal gradient for depth',
        '✓ Rounded corners: 8px → 10px (more modern)',
    ], 7.3, 1.6, 5.5, 4.0, 11, spacing=6)

    add_img(slide, SS_DIR / '07_login_page_375.png', 0.5, 5.3, 2.2, 2.0)
    add_text(slide, 'Mobile (375px): Brand panel hidden, form centered', 2.9, 5.8, 4, 0.4, 11, False, GRAY_MID)

    # ========== SLIDE 6: Welcome Screen ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '04  Welcome Screen — Premium Cards', 0.8, 0.4, 10, 0.6, 28, True)

    add_img(slide, SS_DIR / '02_welcome_screen_1440.png', 0.5, 1.1, 7.5, 4.2)

    add_text(slide, 'Card Button Redesign', 8.3, 1.1, 4.5, 0.4, 18, True)
    add_bullets(slide, [
        '✓ Replaced Ant Design Buttons with custom card-style divs',
        '✓ Hover: elevation (translateY -2px) + yellow border',
        '✓ Yellow icon accent for visual hierarchy',
        '✓ Generous padding (16px) and gap for breathing room',
        '✓ Smooth cubic-bezier transitions (0.2s)',
        '✓ Minimum height (72px) for consistent sizing',
    ], 8.3, 1.6, 4.5, 2.0, 11, spacing=5)

    add_text(slide, 'Logo Upgrade', 8.3, 3.8, 4.5, 0.4, 18, True)
    add_bullets(slide, [
        '✓ Gradient: 145deg #FFE500 → #FDD800',
        '✓ Layered shadow: outer (24px) + inner (8px)',
        '✓ Border-radius: 16px → 18px',
        '✓ fadeInUp entrance animation (0.5s)',
    ], 8.3, 4.3, 4.5, 1.5, 11, spacing=5)

    add_img(slide, SS_DIR / '08_welcome_screen_375.png', 0.5, 5.5, 2.2, 1.8)
    add_text(slide, 'Mobile: cards stack full-width (xs={24})', 2.9, 5.9, 4, 0.4, 11, False, GRAY_MID)

    # ========== SLIDE 7: Chat Page ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '05  Chat Page — Input & Streaming', 0.8, 0.4, 10, 0.6, 28, True)

    add_text(slide, 'Input Area Redesign', 0.8, 1.2, 6, 0.4, 18, True)
    add_bullets(slide, [
        '✓ Removed character count — reduces visual clutter',
        '✓ Floating bar: upward shadow 0 -2px 10px rgba(0,0,0,0.03)',
        '✓ Rounded top corners (radius-lg) for card appearance',
        '✓ Generous bottom padding (24px) for comfortable typing',
        '✓ Send button icon-only when input has content',
        '✓ Fixed height: replaced calc() with flexbox layout',
        '✓ Streaming scroll: instant during streaming, smooth for completed',
        '✓ Only scrolls when user is near bottom (< 100px)',
    ], 0.8, 1.7, 6, 3.0, 11, spacing=5)

    add_text(slide, 'Streaming & Thinking', 7.5, 1.2, 5.5, 0.4, 18, True)
    add_bullets(slide, [
        '✓ Cursor: 2px thin yellow line with subtle glow',
        '✓ Smooth blink (0.8s ease-in-out)',
        '✓ Thinking indicator: 3 animated yellow dots',
        '✓ dotBounce animation with staggered delays',
        '✓ Replaced Ant Design Spin with custom CSS dots',
    ], 7.5, 1.7, 5.5, 2.5, 11, spacing=5)

    add_text(slide, 'Message Bubble Fixes', 7.5, 4.0, 5.5, 0.4, 18, True)
    add_bullets(slide, [
        '✓ User messages: overflow-wrap + word-break for long URLs',
        '✓ Markdown: overflow wrapping + pre/code scroll',
        '✓ Citations: overflow:hidden + ellipsis for long titles',
        '✓ Assistant bubbles: hover elevation (shadow-md + -1px)',
        '✓ Message entrance: fadeInUp (0.3s, staggered)',
    ], 7.5, 4.5, 5.5, 2.5, 11, spacing=5)

    # ========== SLIDE 8: Profile & Dark Mode ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '06  Profile Page & Dark Mode', 0.8, 0.4, 10, 0.6, 28, True)

    add_img(slide, SS_DIR / '03_profile_page_1440.png', 0.5, 1.1, 6.0, 3.6)
    add_text(slide, 'Light Mode', 0.5, 4.8, 6, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_img(slide, SS_DIR / '04_profile_dark_1440.png', 6.8, 1.1, 6.0, 3.6)
    add_text(slide, 'Dark Mode', 6.8, 4.8, 6, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_text(slide, 'Profile Page Improvements', 0.8, 5.3, 11.5, 0.4, 16, True)
    add_bullets(slide, [
        '✓ Responsive maxWidth: min(680px, 100%) with centered layout',
        '✓ Dark mode toggle button in header (Sun/Moon icon)',
        '✓ Ant Design tokens: borderRadius 10px, controlHeight 42px',
        '✓ Segmented control for theme selection (Light/Dark/System)',
        '✓ Dark mode applies to all components via CSS variables',
        '✓ Menu items: rounded pills (10px) with yellow left indicator bar',
    ], 0.8, 5.7, 12, 1.5, 11, spacing=4)

    # ========== SLIDE 9: History & Knowledge Base ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '07  History & Knowledge Base Pages', 0.8, 0.4, 10, 0.6, 28, True)

    add_img(slide, SS_DIR / '05_history_page_1440.png', 0.5, 1.1, 6.0, 3.6)
    add_text(slide, 'History Page', 0.5, 4.8, 6, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_img(slide, SS_DIR / '06_knowledge_base_1440.png', 6.8, 1.1, 6.0, 3.6)
    add_text(slide, 'Knowledge Base Page', 6.8, 4.8, 6, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_text(slide, 'Knowledge Base Improvements', 0.8, 5.3, 11.5, 0.4, 16, True)
    add_bullets(slide, [
        '✓ Table horizontal scroll: scroll={{ x: "max-content" }}',
        '✓ Upload validation: file type (.pdf/.doc/.docx/.txt/.csv/.xlsx/.pptx) + 50MB limit',
        '✓ Delete confirmation modal: Modal.confirm() with danger type',
        '✓ Custom empty state with helpful message',
        '✓ Refined table: borderRadius 10px, headerBorderRadius 10px',
    ], 0.8, 5.7, 12, 1.5, 11, spacing=4)

    # ========== SLIDE 10: Responsive ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '08  Responsive Design — All Breakpoints', 0.8, 0.4, 10, 0.6, 28, True)

    add_img(slide, SS_DIR / '10_welcome_screen_768.png', 0.5, 1.1, 4.5, 5.5)
    add_text(slide, 'Tablet (768px)', 0.5, 6.7, 4.5, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_img(slide, SS_DIR / '02_welcome_screen_1440.png', 5.2, 1.1, 4.0, 2.5)
    add_text(slide, 'Desktop (1440px)', 5.2, 3.7, 4, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_img(slide, SS_DIR / '08_welcome_screen_375.png', 5.2, 4.1, 2.5, 2.8)
    add_text(slide, 'Mobile (375px)', 5.2, 7.0, 2.5, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_img(slide, SS_DIR / '09_profile_page_375.png', 8.0, 4.1, 2.5, 2.8)
    add_text(slide, 'Profile (375px)', 8.0, 7.0, 2.5, 0.3, 12, True, EY_BLACK, PP_ALIGN.CENTER)

    add_text(slide, 'Breakpoint System', 9.5, 1.1, 3.5, 0.4, 16, True)
    add_bullets(slide, [
        '• 375px: Mobile — sider collapses, single column',
        '• 576px: Small — 2-column quick actions',
        '• 768px: Tablet — sider collapses to 64px icons',
        '• 992px: Desktop — full sider expanded',
        '• 1200px: Large — optimal spacing',
        '• 1440px+: Wide — centered max-width containers',
    ], 9.5, 1.6, 3.5, 3.0, 10, spacing=6)

    add_bullets(slide, [
        '• Login: brand panel hides < 800px, stacks vertically',
        '• Sider: collapses at "md" (768px), was "lg" (992px)',
        '• Welcome: xs={24} sm={12} md={8} responsive columns',
        '• Knowledge Base: scroll={{ x: "max-content" }}',
        '• Profile: maxWidth: min(680px, 100%)',
    ], 9.5, 4.0, 3.5, 3.0, 9, GRAY_MID, 5)

    # ========== SLIDE 11: Summary Metrics ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_bar(slide, 0, 0, 13.33)
    add_text(slide, '09  Summary & Metrics', 0.8, 0.4, 10, 0.6, 28, True)

    metrics = [
        ('10', 'Files Modified', 'All frontend pages and shared styles', EY_YELLOW),
        ('12', 'Critical Bugs Fixed', 'Overflow, responsive, validation', GREEN),
        ('20+', 'UI Improvements', 'Tokens, polish, animations', RGBColor(0x18, 0x90, 0xFF)),
        ('6', 'New Animations', 'fadeInUp, fadeIn, dotBounce, shimmer, pulseGlow, blink', RGBColor(0x72, 0x2E, 0xD1)),
        ('5', 'Breakpoints', '375px to 1440px+ responsive', RGBColor(0xFA, 0x8C, 0x16)),
        ('100%', 'Pages Covered', 'Login, Chat, Profile, History, KB', RGBColor(0x52, 0xC4, 0x1A)),
    ]
    for i, (num, label, desc, accent) in enumerate(metrics):
        x = 0.8 + (i % 3) * 4.1
        y = 1.2 + (i // 3) * 1.8
        add_card(slide, x, y, 3.8, 1.5, label, num, desc, accent)

    add_text(slide, 'Modified Files', 0.8, 5.2, 6, 0.4, 16, True)
    add_bullets(slide, [
        '1. index.html — Plus Jakarta Sans Google Font',
        '2. globals.css — Tokens, shadows, animations, responsive, scrollbar',
        '3. ChatPage.tsx — Height fix, streaming scroll, input polish, thinking dots',
        '4. MessageBubble.tsx — Overflow fixes, cursor, hover, entrance animation',
        '5. WelcomeScreen.tsx — Premium card buttons, logo upgrade',
        '6. AppLayout.tsx — Sider breakpoint, header, email truncation, dark toggle',
        '7. useTheme.ts — Token refinement (borderRadius 10, controlHeight 42)',
        '8. App.tsx — Login responsive, brand panel polish, submit hover',
        '9. KnowledgeBasePage.tsx — Table scroll, upload validation, delete confirm, empty state',
        '10. ProfilePage.tsx — Responsive maxWidth',
    ], 0.8, 5.6, 12, 1.8, 10, spacing=3)

    # ========== SLIDE 12: Closing ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = EY_BLACK
    add_bar(slide, 0, 0, 13.33, 0.08)

    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(1.8), Inches(2.5), Inches(1.0))
    logo.fill.solid()
    logo.fill.fore_color.rgb = EY_YELLOW
    logo.line.fill.background()
    tf = logo.text_frame
    p = tf.paragraphs[0]
    p.text = 'EY'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = EY_BLACK
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, 'UI Premium Optimization Complete', 1.5, 3.2, 10.33, 0.6, 30, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, 'All pages optimized with premium typography, refined design tokens,\nresponsive layouts, and subtle micro-interactions.',
             2.5, 3.9, 8.33, 0.8, 14, False, RGBColor(0xA6, 0xA6, 0xA6), PP_ALIGN.CENTER)

    add_text(slide, f'Report generated: {datetime.datetime.now().strftime("%Y-%m-%d %H:%M")}  |  Screenshots: 10',
             3.5, 5.0, 6.33, 0.4, 12, False, GRAY_MID, PP_ALIGN.CENTER)

    # Save
    prs.save(str(OUTPUT_FILE))
    print(f'[OK] Report saved: {OUTPUT_FILE}')
    print(f'    Slides: {len(prs.slides)}')
    print(f'    Size: {OUTPUT_FILE.stat().st_size / 1024:.0f} KB')


if __name__ == '__main__':
    main()
