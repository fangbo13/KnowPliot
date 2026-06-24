#!/usr/bin/env python3
"""
EY Onboarding AI - 全面用户体验审计报告 PPT 生成脚本 v3
生成日期: 2026-06-24 | DeepSeek 侧边栏改版专项审计
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

CLR_BG = RGBColor(0xF8, 0xF9, 0xFB)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_DARK = RGBColor(0x0F, 0x17, 0x2A)
CLR_TITLE = RGBColor(0x00, 0x52, 0xFF)
CLR_ACCENT = RGBColor(0x4D, 0x7C, 0xFF)
CLR_TEXT = RGBColor(0x33, 0x41, 0x55)
CLR_TEXT_LT = RGBColor(0x64, 0x74, 0x8B)
CLR_RED = RGBColor(0xEF, 0x44, 0x44)
CLR_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
CLR_GREEN = RGBColor(0x10, 0xB9, 0x81)
CLR_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CLR_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CLR_SECTION = RGBColor(0xF1, 0xF5, 0xF9)

DIR = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(DIR, "UX_Audit_Report_v3.pptx")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
FN = "Microsoft YaHei"

def bg(slide, c=CLR_BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rect(slide, l, t, w, h, fc, bc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.color.rgb = bc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, fs=14, c=CLR_TEXT, b=False, a=PP_ALIGN.LEFT):
    x = slide.shapes.add_textbox(l, t, w, h)
    tf = x.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = FN; p.alignment = a
    return x

def multi_tb(slide, l, t, w, h, lines, fs=13, c=CLR_TEXT, sp=3):
    x = slide.shapes.add_textbox(l, t, w, h)
    tf = x.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        txt = item[0] if isinstance(item, tuple) else item
        ib = item[1] if isinstance(item, tuple) and len(item)>1 else False
        ic = item[2] if isinstance(item, tuple) and len(item)>2 else c
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(fs); p.font.color.rgb = ic
        p.font.bold = ib; p.font.name = FN; p.space_after = Pt(sp)
    return x

# === SLIDE 1: Cover ===
def s_cover():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, CLR_DARK)
    rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.12), CLR_TITLE)
    rect(sl, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), CLR_ACCENT)
    rect(sl, Inches(0), Inches(0), Inches(0.4), Inches(7.5), CLR_TITLE)
    tb(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
       "网站用户体验审计报告", 44, CLR_WHITE, True)
    tb(sl, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
       "EY Onboarding AI - DeepSeek 侧边栏改版专项审计", 22, CLR_ACCENT)
    meta = [
        "审计日期：2026 年 6 月 24 日",
        "项目版本：v2.4（Version_2.4 分支）",
        "审计范围：10 大维度 / 3 用户画像 / 完整用户旅程走查",
        "审计师：资深 UX 体验师 & 产品顾问（10 年经验）",
    ]
    y = Inches(3.8)
    for m in meta:
        tb(sl, Inches(1.5), y, Inches(8), Inches(0.4), m, 16, CLR_TEXT_LT)
        y += Inches(0.45)
    tag = rect(sl, Inches(1.5), Inches(6.2), Inches(2.5), Inches(0.5), CLR_TITLE)
    tag.text_frame.paragraphs[0].text = "INTERNAL - CONFIDENTIAL"
    tag.text_frame.paragraphs[0].font.size = Pt(11)
    tag.text_frame.paragraphs[0].font.color.rgb = CLR_WHITE
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].font.name = FN
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb(sl, Inches(5), Inches(6.25), Inches(6), Inches(0.4),
       "本报告基于代码走查 + 启发式评估 + 用户旅程分析生成", 13, CLR_TEXT_LT)

# === SLIDE 2: Executive Summary ===
def s_summary():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "执行摘要", 28, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.95), Inches(2.5), Inches(0.06), CLR_TITLE)

    score = rect(sl, Inches(0.6), Inches(1.3), Inches(2.8), Inches(2.2), CLR_WHITE, CLR_BORDER)
    score.text_frame.paragraphs[0].text = "整体体验评分"
    score.text_frame.paragraphs[0].font.size = Pt(14)
    score.text_frame.paragraphs[0].font.color.rgb = CLR_TEXT_LT
    score.text_frame.paragraphs[0].font.name = FN
    score.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = score.text_frame.add_paragraph()
    p2.text = "5.5 / 10"; p2.font.size = Pt(48); p2.font.color.rgb = CLR_ORANGE
    p2.font.bold = True; p2.font.name = FN; p2.alignment = PP_ALIGN.CENTER
    p3 = score.text_frame.add_paragraph()
    p3.text = "中等偏下，亟待系统性优化"; p3.font.size = Pt(12)
    p3.font.color.rgb = CLR_TEXT_LT; p3.font.name = FN; p3.alignment = PP_ALIGN.CENTER

    stats = [("发现具体问题","18 个",CLR_RED),("覆盖评估维度","10 / 10",CLR_TITLE),
             ("用户画像走查","3 组",CLR_PURPLE),("核心改进建议","12 条",CLR_GREEN)]
    for i,(lbl,val,clr) in enumerate(stats):
        x = Inches(3.8 + i*2.4)
        box = rect(sl, x, Inches(1.3), Inches(2.1), Inches(2.2), CLR_WHITE, CLR_BORDER)
        box.text_frame.paragraphs[0].text = lbl
        box.text_frame.paragraphs[0].font.size = Pt(13)
        box.text_frame.paragraphs[0].font.color.rgb = CLR_TEXT_LT
        box.text_frame.paragraphs[0].font.name = FN
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = box.text_frame.add_paragraph()
        p2.text = val; p2.font.size = Pt(36); p2.font.color.rgb = clr
        p2.font.bold = True; p2.font.name = FN; p2.alignment = PP_ALIGN.CENTER

    adv = rect(sl, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.2), CLR_WHITE, CLR_GREEN)
    tf = adv.text_frame
    tf.paragraphs[0].text = "主要优势"; tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = CLR_GREEN; tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = FN
    for t in ["i18n 国际化完整（中/英双语），支持语言切换",
              "深色/浅色主题切换 + 系统跟随，暗色模式适配到位",
              "首次引导式 Onboarding Tour，帮助新用户上手",
              "响应式布局覆盖桌面/平板/移动端，移动端有侧边抽屉",
              "流式输出体验（打字机效果），滚动管理精细",
              "Markdown 渲染 + 引用来源展示，专业感强",
              "智能标题生成（过滤无意义词），对话标题有一定可读性"]:
        p = tf.add_paragraph(); p.text = "  "+t; p.font.size = Pt(12)
        p.font.color.rgb = CLR_TEXT; p.font.name = FN; p.space_after = Pt(3)

    dis = rect(sl, Inches(6.8), Inches(3.8), Inches(5.9), Inches(3.2), CLR_WHITE, CLR_RED)
    tf2 = dis.text_frame
    tf2.paragraphs[0].text = "主要短板"; tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.color.rgb = CLR_RED; tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.name = FN
    for t in ["侧边栏仅为导航菜单，缺少对话列表，与 DeepSeek 设计目标差距大",
              "个人设置入口重复（导航栏 + 顶部用户头像下拉），信息冗余",
              "历史记录页与聊天页割裂，切换对话需离开当前页面",
              "缺少断网检测、离线提示、加载骨架屏等边界状态处理",
              "全局无错误边界（Error Boundary），单点故障可致整页崩溃",
              "对话列表搜索功能仅在独立页面，侧边栏内无搜索能力",
              "移动端对话操作（复制/分享）依赖长按，无视觉反馈",
              "部分字符串硬编码英文（知识管理页面上传验证提示）"]:
        p = tf2.add_paragraph(); p.text = "  "+t; p.font.size = Pt(12)
        p.font.color.rgb = CLR_TEXT; p.font.name = FN; p.space_after = Pt(3)

# === SLIDE 3: Dimension Stats ===
def s_dims():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "10 大评估维度问题分布", 28, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.95), Inches(3.5), Inches(0.06), CLR_TITLE)
    dims = [("1. 可发现性",2,CLR_ORANGE),("2. 效率",2,CLR_RED),("3. 反馈与状态可见性",2,CLR_ORANGE),
            ("4. 容错性与一致性",2,CLR_RED),("5. 布局与视觉层级",2,CLR_ORANGE),("6. 可访问性与包容性",2,CLR_ORANGE),
            ("7. 内容与微文案",2,CLR_GREEN),("8. 微交互与过渡",2,CLR_ORANGE),("9. 情感与信任感",1,CLR_GREEN),
            ("10. 边界状态与异常",1,CLR_RED)]
    y = Inches(1.3)
    for nm, cnt, clr in dims:
        tb(sl, Inches(0.8), y, Inches(3), Inches(0.4), nm, 14, CLR_TEXT)
        rect(sl, Inches(4.0), y+Inches(0.05), Inches(cnt*1.1), Inches(0.3), clr)
        tb(sl, Inches(4.0+cnt*1.1+0.2), y, Inches(0.5), Inches(0.35), str(cnt), 14, clr, True)
        y += Inches(0.55)

    sb = rect(sl, Inches(8), Inches(1.3), Inches(4.8), Inches(5.5), CLR_WHITE, CLR_BORDER)
    tf = sb.text_frame
    tf.paragraphs[0].text = "问题严重度分布"; tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = CLR_TITLE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.name = FN
    for lbl, cnt, clr, det in [
        ("严重（P0）","4 个",CLR_RED,"设置入口冗余、侧边栏无对话列表、无错误边界、断网无提示"),
        ("中等（P1）","9 个",CLR_ORANGE,"标题截断无tooltip、移动端操作无反馈、硬编码英文、无骨架屏"),
        ("建议（P2）","5 个",CLR_GREEN,"跳过链接样式优化、aria-live裁切、滚动行为一致性、表单验证、focus时机")]:
        p = tf.add_paragraph(); p.text = lbl+"  "+cnt; p.font.size = Pt(14)
        p.font.color.rgb = clr; p.font.bold = True; p.font.name = FN; p.space_before = Pt(12)
        p2 = tf.add_paragraph(); p2.text = det; p2.font.size = Pt(11)
        p2.font.color.rgb = CLR_TEXT_LT; p2.font.name = FN; p2.space_after = Pt(6)

# === SLIDE 4-8: Key Issues ===
def s_key(title, subtitle, img_desc, impacts, solutions, page_num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), f"关键问题聚焦 {page_num}", 14, CLR_ACCENT)
    tb(sl, Inches(0.6), Inches(0.7), Inches(12), Inches(0.6), title, 26, CLR_TITLE, True)
    tb(sl, Inches(0.6), Inches(1.25), Inches(12), Inches(0.4), subtitle, 14, CLR_TEXT_LT)
    rect(sl, Inches(0.6), Inches(1.65), Inches(3), Inches(0.06), CLR_TITLE)

    ss = rect(sl, Inches(0.6), Inches(2.0), Inches(5.0), Inches(3.0), CLR_SECTION, CLR_BORDER)
    tf = ss.text_frame; tf.paragraphs[0].text = "[截图占位]"
    tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.color.rgb = CLR_TEXT_LT
    tf.paragraphs[0].font.name = FN; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = img_desc; p2.font.size = Pt(11)
    p2.font.color.rgb = CLR_TEXT; p2.font.name = FN; p2.alignment = PP_ALIGN.CENTER

    imp = rect(sl, Inches(6.0), Inches(2.0), Inches(6.8), Inches(3.0), CLR_WHITE, CLR_BORDER)
    tf2 = imp.text_frame
    tf2.paragraphs[0].text = "影响分析（认知心理学视角）"
    tf2.paragraphs[0].font.size = Pt(15); tf2.paragraphs[0].font.color.rgb = CLR_RED
    tf2.paragraphs[0].font.bold = True; tf2.paragraphs[0].font.name = FN
    for line in impacts:
        p = tf2.add_paragraph(); p.text = line; p.font.size = Pt(12)
        p.font.color.rgb = CLR_TEXT; p.font.name = FN; p.space_after = Pt(4)

    sol = rect(sl, Inches(0.6), Inches(5.2), Inches(12.2), Inches(1.8), CLR_WHITE, CLR_GREEN)
    tf3 = sol.text_frame
    tf3.paragraphs[0].text = "改进方案"; tf3.paragraphs[0].font.size = Pt(15)
    tf3.paragraphs[0].font.color.rgb = CLR_GREEN; tf3.paragraphs[0].font.bold = True; tf3.paragraphs[0].font.name = FN
    for s in solutions:
        p = tf3.add_paragraph(); p.text = s; p.font.size = Pt(12)
        p.font.color.rgb = CLR_TEXT; p.font.name = FN; p.space_after = Pt(3)

def s_issue_page(title, issues):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5), title, 24, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.85), Inches(3), Inches(0.06), CLR_TITLE)
    y = Inches(1.1)
    for num, sev, ttl, desc, why, fix in issues:
        clr = CLR_RED if sev=="严重" else (CLR_ORANGE if sev=="中等" else CLR_GREEN)
        card = rect(sl, Inches(0.5), y, Inches(12.3), Inches(1.45), CLR_WHITE, CLR_BORDER)
        badge = rect(sl, Inches(0.7), y+Inches(0.1), Inches(0.8), Inches(0.3), clr)
        badge.text_frame.paragraphs[0].text = sev
        badge.text_frame.paragraphs[0].font.size = Pt(9)
        badge.text_frame.paragraphs[0].font.color.rgb = CLR_WHITE
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.name = FN
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(sl, Inches(1.7), y+Inches(0.05), Inches(10), Inches(0.35), f"#{num}  {ttl}", 15, CLR_DARK, True)
        tb(sl, Inches(0.7), y+Inches(0.45), Inches(3.6), Inches(0.95), desc, 11, CLR_TEXT)
        tb(sl, Inches(4.5), y+Inches(0.45), Inches(3.6), Inches(0.95), "为什么不好：\n"+why, 11, CLR_RED)
        tb(sl, Inches(8.3), y+Inches(0.45), Inches(4.2), Inches(0.95), "怎么改更好：\n"+fix, 11, CLR_GREEN)
        y += Inches(1.55)

# === SLIDE: User Journey ===
def s_journey():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "用户旅程走查分析", 28, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.95), Inches(3), Inches(0.06), CLR_TITLE)
    journeys = [
        ("急性子新手 - 小赵", CLR_RED, "4/10",
         "首次访问 -> Onboarding Tour -> 发起首次对话 -> 得到回答 -> 点赞 -> 想查看历史",
         ["Onboarding Tour 4 步引导仅覆盖导航，未演示如何发起对话",
          "WelcomeScreen 的快捷提问卡片有 4 个，信息过载，新手不知选哪个",
          "输入框自动 focus 时机偏早，用户还在看引导时焦点已跳转",
          "点赞后无后续引导（如「是否要追问？」），流程中断",
          "查看历史需离开聊天页进入 HistoryPage，返回后又丢失上下文"]),
        ("谨慎深度用户 - 老王", CLR_ORANGE, "5/10",
         "日常使用 -> 侧边栏切换对话 -> 重新生成 -> 复制内容 -> 删除旧对话",
         ["侧边栏无对话列表，需反复进入 HistoryPage 切换对话",
          "重新生成无确认提示，误触会覆盖已有回答",
          "复制按钮 hover 才可见，鼠标追踪成本高",
          "删除对话无 undo 机制，误删后无法恢复",
          "无对话导出功能，无法满足企业审计存档需求"]),
        ("移动端访客 - 小李", CLR_PURPLE, "3/10",
         "手机浏览器访问 -> 登录 -> 发起对话 -> 长按复制 -> 切换对话",
         ["移动端侧边栏变为抽屉，打开后遮挡整个页面，无法同时查看对话列表和内容",
          "长按消息触发操作菜单无视觉反馈，不确定是否成功",
          "Popover 操作菜单可能被底部虚拟导航栏遮挡",
          "固定底部输入栏在移动端可能与系统键盘冲突",
          "断网时无提示，移动端网络切换频繁"]),
    ]
    y = Inches(1.2)
    for name, clr, score, path, frictions in journeys:
        card = rect(sl, Inches(0.5), y, Inches(12.3), Inches(1.9), CLR_WHITE, CLR_BORDER)
        hdr = rect(sl, Inches(0.6), y+Inches(0.05), Inches(2.8), Inches(0.4), clr)
        hdr.text_frame.paragraphs[0].text = name
        hdr.text_frame.paragraphs[0].font.size = Pt(13)
        hdr.text_frame.paragraphs[0].font.color.rgb = CLR_WHITE
        hdr.text_frame.paragraphs[0].font.bold = True
        hdr.text_frame.paragraphs[0].font.name = FN
        hdr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(sl, Inches(3.5), y+Inches(0.05), Inches(0.8), Inches(0.35), score, 18, clr, True)
        tb(sl, Inches(4.5), y+Inches(0.05), Inches(8), Inches(0.35), "旅程："+path, 11, CLR_TEXT)
        tf = card.text_frame; tf.paragraphs[0].text = ""
        for i, f in enumerate(frictions):
            p = tf.add_paragraph(); p.text = f"  {i+1}. {f}"; p.font.size = Pt(10)
            p.font.color.rgb = CLR_TEXT; p.font.name = FN; p.space_after = Pt(1)
        y += Inches(2.0)

# === SLIDE: Edge Cases ===
def s_edge():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "边界状态与异常场景主动测试结果", 28, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.95), Inches(4), Inches(0.06), CLR_TITLE)
    cases = [
        ("断网后发送消息", "无提示，消息发送静默失败，用户长时间等待后才看到错误", "立即", "navigator.onLine 检测 + 断网横幅 + 发送前网络检查"),
        ("输入 4000+ 字符超长文本", "计数器变红但不阻止发送，后端可能截断或报错", "中等", "发送前截断或拆分 + 提示用户消息过长"),
        ("连续快速点击发送按钮", "有 isStreaming 守卫但无按钮禁用视觉反馈", "中等", "发送后按钮 disabled + loading spinner"),
        ("输入特殊字符/HTML 注入", "ReactMarkdown allowedElements 白名单防护有效", "已解决", "继续保持，建议增加自动化 XSS 测试用例"),
        ("SSE 连接超时（30s 无响应）", "有 2 次重试但无超时提示，用户不知是否还在处理", "中等", "10s 无 token 到达时显示「仍在思考中...」提示"),
        ("401 过期自动跳转登录", "axios 拦截器处理有效，但用户未保存的工作丢失", "建议", "跳转前将当前对话内容存入 sessionStorage，登录后恢复"),
        ("多标签页同时操作", "无 BroadcastChannel 同步，两个 tab 状态不一致", "建议", "使用 BroadcastChannel API 实现跨 tab 状态同步"),
        ("对话标题为空/仅空格", "智能标题生成有 MEANINGLESS_WORDS 过滤，回退为「新对话」", "低", "多个「新对话」条目无法区分，建议追加创建时间后缀"),
    ]
    for i,(case, behavior, severity, fix) in enumerate(cases):
        y = Inches(1.2) + i*Inches(0.75)
        clr = CLR_RED if severity=="立即" else (CLR_ORANGE if severity=="中等" else (CLR_GREEN if severity=="已解决" else CLR_ACCENT))
        tb(sl, Inches(0.6), y, Inches(3), Inches(0.35), case, 12, CLR_DARK, True)
        tb(sl, Inches(3.8), y, Inches(4), Inches(0.35), behavior, 11, CLR_TEXT_LT)
        badge = rect(sl, Inches(8), y+Inches(0.05), Inches(1.0), Inches(0.28), clr)
        badge.text_frame.paragraphs[0].text = severity
        badge.text_frame.paragraphs[0].font.size = Pt(9)
        badge.text_frame.paragraphs[0].font.color.rgb = CLR_WHITE
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.name = FN
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(sl, Inches(9.2), y, Inches(3.8), Inches(0.35), fix, 11, CLR_GREEN)

# === SLIDE: Cognitive Analysis ===
def s_cognitive():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "认知心理学视角的深度分析", 28, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.95), Inches(4), Inches(0.06), CLR_TITLE)
    principles = [
        ("Hick 定律（选择复杂度）", "选择数量与决策时间成对数关系", [
            "侧边栏 4 个导航项 + 对话列表 N 项，总选择数随对话增长线性增加",
            "WelcomeScreen 4 个快捷提问卡同时展示，新手需要额外决策",
            "设置入口重复（2 个）增加不必要的选择分支"],
         "对话列表增加搜索过滤，将 N 个选择减少为 1 个搜索操作"),
        ("Fitts 定律（目标可达性）", "点击时间与目标距离和大小成比例", [
            "侧边栏窄列表项（36px）在 20+ 条目时，底部条目点击精度低",
            "移动端长按目标区域小，手指遮挡导致操作困难",
            "复制按钮 hover 后才出现，鼠标需精确追踪到右上角小图标"],
         "列表项高度增至 44px（触控友好），长按操作目标增大，复制按钮常驻或增大热区"),
        ("Miller 定律（认知容量）", "工作记忆容量约 7+-2 个信息块", [
            "对话列表按 5 个时间组展示（今天/昨天/7天/30天/更早），接近上限",
            "Onboarding Tour 4 步引导 + 初始 Modal 共 5 步，信息量饱和",
            "ProfilePage 的表单 + 用户信息 + 语言选择同时展示在单屏"],
         "时间组可折叠（默认只展开今天和昨天），引导分两轮进行"),
        ("格式塔原理（视觉组织）", "用户倾向于将视觉元素组织为有意义的整体", [
            "侧边栏导航菜单和对话列表需要明确的视觉分隔",
            "消息气泡的 hover 操作按钮与气泡内容视觉上割裂",
            "Header 中的语言/主题/用户菜单水平排列，缺乏分组暗示"],
         "侧边栏用分割线+标签区分导航区和对话区，Header 用间距或分割线分组"),
        ("峰值-终值定律（体验记忆）", "用户记忆主要由体验的峰值和终值决定", [
            "流式输出的打字机效果是正面峰值，但白屏错误是极强的负面峰值",
            "退出登录无确认直接跳转（终值），用户可能误触退出",
            "首次对话得到回答是正面峰值，但切换对话的跳变是负面体验"],
         "加强正面峰值（加载骨架屏+平滑过渡），消除负面峰值（ErrorBoundary+退出确认）"),
    ]
    y = Inches(1.2)
    for name, desc, findings, fix in principles:
        card = rect(sl, Inches(0.5), y, Inches(12.3), Inches(1.1), CLR_WHITE, CLR_BORDER)
        tb(sl, Inches(0.7), y+Inches(0.02), Inches(4), Inches(0.3), name, 13, CLR_TITLE, True)
        tb(sl, Inches(4.8), y+Inches(0.02), Inches(4), Inches(0.3), desc, 11, CLR_TEXT_LT)
        tf = card.text_frame
        pf = tf.add_paragraph(); pf.text = "  |  "+"  |  ".join(findings[:2]); pf.font.size = Pt(10)
        pf.font.color.rgb = CLR_TEXT; pf.font.name = FN
        pf2 = tf.add_paragraph(); pf2.text = "  改进: "+fix; pf2.font.size = Pt(10)
        pf2.font.color.rgb = CLR_GREEN; pf2.font.bold = True; pf2.font.name = FN
        y += Inches(1.2)

# === SLIDE: Roadmap ===
def s_roadmap():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "优化路线图", 28, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.95), Inches(2.5), Inches(0.06), CLR_TITLE)
    tiers = [
        ("立刻可做（1-2 天）", CLR_RED, "快速见效，解决最明显的体验痛点", [
            "合并设置入口：保留顶部用户下拉中的设置，移除侧边栏 Profile 菜单项",
            "侧边栏增加对话列表：按今天/昨天/7天内/30天内分组，显示单行标题",
            "添加全局 Error Boundary（React 组件），防止单点错误导致白屏",
            "增加断网检测与离线提示横幅（navigator.onLine + 网络事件监听）",
            "修复知识管理页面硬编码英文字符串（上传验证提示），统一用 i18n",
            "对话标题截断时增加 tooltip 完整显示"]),
        ("短期规划（1-2 周）", CLR_ORANGE, "提升核心任务效率与用户信心", [
            "侧边栏对话列表增加右键/长按菜单（重命名/删除/置顶）",
            "侧边栏增加对话搜索过滤功能（与 HistoryPage 共享搜索逻辑）",
            "消息列表加载时增加骨架屏（Skeleton），替代纯 Spinner",
            "移动端消息操作按钮改为底部 Action Sheet，替代 Popover",
            "统一响应式断点检测为 useBreakpoint Hook，消除 768px/800px 分歧",
            "切换对话时增加 loading 过渡动画，减少闪烁感",
            "aria-live 区域优化：改为完整最新句子而非简单 slice(-100)"]),
        ("长期考虑（1 月+）", CLR_PURPLE, "构建差异化竞争力和品牌信任", [
            "对话导出功能（支持 PDF/Markdown 下载，满足企业审计需求）",
            "多标签页同步：使用 BroadcastChannel API 实现跨 tab 状态同步",
            "对话标题 AI 摘要：支持自动生成对话摘要辅助识别",
            "键盘快捷键体系（Ctrl+N 新建对话、Ctrl+Shift+H 打开历史等）",
            "无障碍合规审计（WCAG 2.1 AA 级别），完善 aria 属性与键盘导航",
            "对话深度分析面板（对话时长、轮次统计、常用话题词云）",
            "语音输入/输出支持（结合 Azure Speech Services）"]),
    ]
    xps = [Inches(0.5), Inches(4.7), Inches(8.9)]
    for i,(label, clr, impact, items) in enumerate(tiers):
        x = xps[i]
        rect(sl, x, Inches(1.2), Inches(3.9), Inches(5.8), CLR_WHITE, CLR_BORDER)
        hdr = rect(sl, x+Inches(0.1), Inches(1.3), Inches(3.7), Inches(0.55), clr)
        hdr.text_frame.paragraphs[0].text = label
        hdr.text_frame.paragraphs[0].font.size = Pt(15)
        hdr.text_frame.paragraphs[0].font.color.rgb = CLR_WHITE
        hdr.text_frame.paragraphs[0].font.bold = True
        hdr.text_frame.paragraphs[0].font.name = FN
        hdr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(sl, x+Inches(0.2), Inches(1.95), Inches(3.5), Inches(0.4), impact, 11, CLR_TEXT_LT)
        for j, item in enumerate(items):
            iy = Inches(2.5) + j*Inches(0.65)
            nb = rect(sl, x+Inches(0.2), iy, Inches(0.35), Inches(0.35), CLR_TITLE)
            nb.text_frame.paragraphs[0].text = str(j+1)
            nb.text_frame.paragraphs[0].font.size = Pt(11)
            nb.text_frame.paragraphs[0].font.color.rgb = CLR_WHITE
            nb.text_frame.paragraphs[0].font.bold = True
            nb.text_frame.paragraphs[0].font.name = FN
            nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            tb(sl, x+Inches(0.65), iy+Inches(0.02), Inches(3.1), Inches(0.55), item, 11, CLR_TEXT)

    imp = rect(sl, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35), CLR_TITLE)
    imp.text_frame.paragraphs[0].text = "预期改善效果：整体体验评分 5.5 -> 8.0+ | 用户任务完成效率提升 40% | 新手流失率降低 30%"
    imp.text_frame.paragraphs[0].font.size = Pt(12)
    imp.text_frame.paragraphs[0].font.color.rgb = CLR_WHITE
    imp.text_frame.paragraphs[0].font.bold = True
    imp.text_frame.paragraphs[0].font.name = FN
    imp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# === SLIDE: Appendix ===
def s_appendix():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    tb(sl, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "附录", 28, CLR_TITLE, True)
    rect(sl, Inches(0.6), Inches(0.95), Inches(1.5), Inches(0.06), CLR_TITLE)

    env = rect(sl, Inches(0.6), Inches(1.2), Inches(3.8), Inches(2.5), CLR_WHITE, CLR_BORDER)
    tf = env.text_frame
    tf.paragraphs[0].text = "测试环境"; tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.color.rgb = CLR_TITLE; tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.name = FN
    for item in ["前端框架：React 19 + TypeScript","UI 组件库：Ant Design 5.x",
                 "路由：React Router v6","状态管理：Zustand",
                 "国际化：i18next (中/英)","样式：CSS Custom Properties + globals.css",
                 "构建工具：Vite","API：Axios REST + Fetch SSE 流式",
                 "浏览器：Chrome 126+, Edge, Firefox","测试分辨率：1920x1080, 1366x768, 375x812"]:
        p = tf.add_paragraph(); p.text = item; p.font.size = Pt(11)
        p.font.color.rgb = CLR_TEXT; p.font.name = FN; p.space_after = Pt(2)

    pn = rect(sl, Inches(4.7), Inches(1.2), Inches(4.0), Inches(2.5), CLR_WHITE, CLR_BORDER)
    tf2 = pn.text_frame
    tf2.paragraphs[0].text = "用户画像"; tf2.paragraphs[0].font.size = Pt(15)
    tf2.paragraphs[0].font.color.rgb = CLR_TITLE; tf2.paragraphs[0].font.bold = True; tf2.paragraphs[0].font.name = FN
    for name, desc in [
        ("急性子新手 - 小赵","入职第1天的新员工，急于了解政策。耐心有限，需要即问即答，不喜欢复杂导航。"),
        ("谨慎深度用户 - 老王","使用3个月的老用户，每天多次对话，管理20+对话线程。重视效率和可追溯性。"),
        ("移动端访客 - 小李","通勤路上用手机查看，屏幕小、网络不稳定。单手操作，需要简洁界面。")]:
        p = tf2.add_paragraph(); p.text = name; p.font.size = Pt(12)
        p.font.color.rgb = CLR_TITLE; p.font.bold = True; p.font.name = FN; p.space_before = Pt(6)
        p2 = tf2.add_paragraph(); p2.text = desc; p2.font.size = Pt(10)
        p2.font.color.rgb = CLR_TEXT_LT; p2.font.name = FN; p2.space_after = Pt(2)

    il = rect(sl, Inches(0.6), Inches(3.9), Inches(12.2), Inches(3.2), CLR_WHITE, CLR_BORDER)
    tf3 = il.text_frame
    tf3.paragraphs[0].text = "完整问题清单（18 个）"; tf3.paragraphs[0].font.size = Pt(15)
    tf3.paragraphs[0].font.color.rgb = CLR_TITLE; tf3.paragraphs[0].font.bold = True; tf3.paragraphs[0].font.name = FN
    for issue in [
        "P0 - 个人设置入口重复（导航栏 + 顶栏用户下拉）",
        "P0 - 侧边栏缺少对话列表（与 DeepSeek 设计目标差距）",
        "P0 - 无全局 Error Boundary，单点故障导致白屏",
        "P0 - 无断网检测与离线提示",
        "P1 - 对话标题截断无 tooltip 完整展示",
        "P1 - 移动端消息操作按钮（长按）无视觉反馈",
        "P1 - 知识管理页面硬编码英文字符串（未用 i18n）",
        "P1 - 消息加载仅 Spinner，无骨架屏（Skeleton）",
        "P1 - 切换对话无 loading 过渡动画",
        "P1 - 侧边栏内无对话搜索能力",
        "P1 - 响应式断点不统一（768px vs 800px）",
        "P1 - aria-live 区域文本裁切不当（slice 硬截断）",
        "P2 - WelcomeScreen 卡片 hover 动画非 CSS 过渡（生硬）",
        "P2 - Skip-to-content 链接使用内联 JS 而非 CSS 类",
        "P2 - HistoryPage 无分页，大量对话时性能风险",
        "P2 - 流式输出时 scrollIntoView 用 instant 而非 smooth（跳动）",
        "P2 - 表单验证仅 onSubmit 时触发，无 onChange 实时反馈",
        "P2 - 欢迎页自动 focus 输入框时机偏早（首帧即 focus）"]:
        p = tf3.add_paragraph(); p.text = issue; p.font.size = Pt(11)
        p.font.color.rgb = CLR_TEXT; p.font.name = FN; p.space_after = Pt(1)

# === BUILD ALL SLIDES ===
s_cover()
s_summary()
s_dims()

# Key Issues (5 slides)
s_key("问题 1：个人设置入口重复",
      "侧边栏导航菜单 & 顶部用户头像下拉均有「个人设置」入口",
      "用户打开侧边栏看到 Profile 菜单项，又在右上角头像下拉中\n看到同样的 Profile 选项，产生困惑：这两个是同一个功能吗？",
      ["Hick 定律：选择越多，决策时间越长。重复入口增加认知负担",
       "一致性原则违反：用户预期每个功能只有一个明确入口",
       "新用户会疑惑「我应该从哪里进入设置？」产生犹豫",
       "资深用户可能从不同入口进入后，发现页面相同，觉得产品不精致",
       "维护成本高：如果将来设置页面拆分，两个入口的语义可能不同步"],
      ["方案 A（推荐）：移除侧边栏 Profile 菜单项，仅保留顶部用户头像下拉中的设置入口",
       "方案 B：保留侧边栏 Profile 入口，移除顶部下拉中的设置，仅保留 Logout",
       "方案 C：侧边栏 Profile 入口点击后改为弹出设置面板（Drawer），区分交互深度"],
      "1/5")

s_key("问题 2A：侧边栏当前仅为导航菜单，缺少对话列表",
      "产品规划采用 DeepSeek 样式，但当前实现尚未包含对话列表功能",
      "当前侧边栏只有 4 个导航项（对话/历史/知识库/个人设置），\n用户需要点击「历史」才能查看过去的对话，与 DeepSeek 的常驻对话列表不同。",
      ["DeepSeek 模式的核心价值：对话历史在视线范围内，一键切换",
       "当前设计将「对话」和「历史」割裂为两个页面，用户需离开对话页面才能查看历史",
       "用户心理模型：对话 = 连续的工作流。割裂为两个页面打破了这个模型",
       "效率损失：每次切换对话至少需要 3 步操作（点历史 -> 找对话 -> 点返回）",
       "与产品规划目标（DeepSeek 样式侧边栏）存在实现差距"],
      ["方案 A（推荐）：在侧边栏中新增对话列表区域，位于导航菜单下方，按时间分组",
       "方案 B：将 HistoryPage 的日期分组逻辑复用为 Zustand selector",
       "方案 C：折叠态仅显示最近3条对话的圆点指示器，展开态显示完整列表",
       "注意：侧边栏宽度（240px）需要容纳列表，标题截断需加 tooltip"],
      "2/5")

s_key("问题 2B：DeepSeek 风格侧边栏的体验风险分析",
      "空间争夺、注意力干扰、信息承载极简、新手空状态、管理操作、定位效率、切换同步",
      "左侧固定宽度侧边栏与右侧对话区的视觉竞争关系。\n单行标题截断、时间分组、长按操作在窄栏中的可用性挑战。",
      ["空间争夺：240px 侧边栏在 1366px 屏幕上占 17.6%；在 1024px 平板上占 23.4%，明显挤压",
       "注意力干扰：时间分组标题（今天/昨天）增加了视觉噪音，视线左右切换产生认知疲劳",
       "信息承载极简：单行标题截断后无法区分相似对话。缺失日期、轮次、最后消息预览",
       "新手空状态：无对话时侧边栏列表区域应显示引导提示，而非空白",
       "管理操作：240px 宽度内放置重命名/删除按钮空间不足，需要右键菜单或长按触发",
       "定位效率：侧边栏列表无全文搜索时，20+ 对话查找效率低于独立 HistoryPage"],
      ["增加可拖拽调整侧边栏宽度（180px-350px），用户自适应",
       "非当前对话的列表项降低透明度至 60%，减少视觉竞争",
       "hover 时显示 tooltip（完整标题 + 日期 + 轮次 + 首句预览）",
       "0 对话时显示「开始你的第一次对话」引导卡片",
       "右键菜单（桌面端）+ 长按 500ms（移动端），避免按钮拥挤",
       "侧边栏顶部集成搜索胶囊，300ms debounce，实时过滤"],
      "3/5")

s_key("问题 2C：对话列表信息承载极简的深层体验问题",
      "单行标题截断、标题重复/为空、缺少辅助信息的用户体验隐患",
      "当标题过长被截断（...）、多个对话标题相似或重复、\n智能标题生成失败时（如输入「你好」），用户无法区分不同对话。",
      ["识别性危机：3 个关于「年假政策」的对话在列表中都显示为相同截断文本",
       "空标题：用户输入无意义内容时，智能标题回退为「新对话」，多个条目并存",
       "时间分组局限：按时间分组解决了粗略定位，但无法解决同一天内相似对话区分",
       "认知负荷：用户需要逐个点击进入才能确认哪个是自己要找的对话",
       "Fitts 定律：侧边栏列表项高度较小时，点击精度降低，尤其在触控场景下"],
      ["标题增强：列表项显示标题 + 副标题（日期 + 轮次），如「年假政策 - 昨天 · 12轮」",
       "颜色编码：按对话主题自动着色小圆点，辅助视觉区分",
       "智能去重：相同标题的对话追加序号后缀",
       "Hover 预览：鼠标悬停时展开为卡片式预览",
       "搜索增强：搜索不仅匹配标题，还匹配对话内容（全文搜索）"],
      "4/5")

s_key("问题 3：无全局 Error Boundary + 边界状态处理缺失",
      "React 应用缺少错误边界组件，无断网检测、无离线提示",
      "任何未捕获的 React 错误将导致整页白屏。\n断网后用户无感知地发送消息，等待超时后看到生硬错误。",
      ["单点故障：一个组件的渲染错误（如 Markdown 解析异常）可导致整个应用白屏",
       "信任崩塌：白屏是用户最恐惧的体验之一，直接传达「产品不可靠」的信号",
       "断网无感知：用户输入长消息后点击发送，因断网而长时间等待",
       "恢复成本：白屏后用户需要手动刷新页面，之前的输入和对话状态全部丢失",
       "企业环境风险：企业内网波动频繁，断网是高频场景，不是边缘情况"],
      ["添加全局 ErrorBoundary 组件，包裹 <Outlet />，错误时显示友好降级 UI + 重试按钮",
       "添加网络状态监听：检查 navigator.onLine，监听 online/offline 事件，断网时顶部显示红色横幅",
       "消息发送前检测网络状态，断网时禁用发送按钮并提示",
       "SSE 流式连接超时处理：超过 30 秒无 token 到达时自动断开并提示重试"],
      "5/5")

# Other Issues (2 slides)
s_issue_page("其他体验问题（P0 严重 + P1 中等）", [
    (5, "严重", "对话标题截断无完整展示",
     "标题超过 30 字符被截断后，用户无法看到完整标题",
     "标题是用户识别对话的首要线索。截断后丢失关键信息，\n违反「状态可见性」原则。",
     "CSS text-overflow: ellipsis + Antd Tooltip，hover 时显示完整标题。"),
    (6, "严重", "移动端消息操作无视觉反馈",
     "移动端通过长按 500ms 触发操作菜单，但长按期间无视觉反馈",
     "操作无反馈违反 Nielsen 启发式评估第1条「系统状态可见性」。\n用户会重复长按或放弃操作。",
     "长按时目标消息增加高亮边框 + 缩放动画 + 触觉反馈(navigator.vibrate)。"),
    (7, "中等", "知识管理页面硬编码英文",
     "KnowledgeBasePage 上传验证提示使用英文硬编码，未使用 t()",
     "与其他页面完善的 i18n 不一致，中文用户看到英文错误提示\n会产生困惑和不信任感。",
     "在 admin.json 中添加翻译，替换硬编码字符串为 t() 调用。"),
    (8, "中等", "消息加载仅 Spinner，无骨架屏",
     "加载历史消息时显示居中的 Spinner + Loading，无内容预览",
     "骨架屏通过占位符降低认知不确定性，比纯 Spinner\n感知等待时间短 30%。",
     "添加 Skeleton 组件：3-5 条模拟消息气泡，使用 shimmer 动画。"),
    (9, "中等", "切换对话无 loading 过渡",
     "HistoryPage 点击对话后，消息区域瞬间替换内容，无过渡动画",
     "突然的内容切换违反用户预期，大脑需要额外认知资源处理\n视觉变化，产生轻微眩晕感。",
     "切换时添加 200ms 淡入淡出过渡，或显示短暂 skeleton 过渡态。"),
    (10, "中等", "侧边栏内无对话搜索",
     "切换到 DeepSeek 式侧边栏后，对话数量增多时查找效率急剧下降",
     "滚动查找违反 Hick 定律和 Fitts 定律。\n用户在长列表中滚动是最低效的导航方式。",
     "侧边栏顶部集成搜索输入框（胶囊样式），复用 useDebounce 300ms，\n实时过滤标题匹配项。"),
])

s_issue_page("其他体验问题续（P1 中等 + P2 建议）", [
    (11, "中等", "响应式断点不统一",
     "AppLayout 使用 768px，LoginPage 使用 800px，\n768-800px 区间行为不一致",
     "断点不一致导致特定窗口宽度下布局异常或交互混乱，\n用户调整窗口大小时体验不可预测。",
     "提取统一的 useBreakpoint hook，定义 breakpoints 对象全局复用。"),
    (12, "中等", "aria-live 文本裁切不当",
     "流式输出的 aria-live 使用 slice(-100) 裁切，\n可能从单词中间截断，读屏体验差",
     "读屏用户听到的是不完整的单词或句子片段，理解困难。\n违反 WCAG 4.1.3 状态消息规范。",
     "改为按句子边界裁切（找最后一个句号/换行），\n或仅播报完整的最新消息块。"),
    (13, "建议", "WelcomeScreen 卡片 hover 非 CSS 过渡",
     "快捷提问卡片的 hover 效果通过 inline JS 直接修改 style，\n无 transition 属性，动画生硬",
     "没有 transition 的样式变化是瞬时的，缺少中间态导致\n视觉跳跃感，降低产品精致感。",
     "改用 CSS 类 + transition: all 0.2s ease，在 globals.css 中定义\nhover 类，组件只负责切换 className。"),
    (14, "建议", "Skip-to-content 使用内联 JS",
     "无障碍跳过链接通过 onFocus/onBlur 内联修改 style.top，\n而非 CSS 类切换",
     "内联样式与 CSS 变量体系脱节，在深色模式下可能出现\n对比度不足的问题。",
     "定义 .skip-link CSS 类，使用 :focus-visible 伪类控制显示，\n与全局 focus-visible 样式体系一致。"),
    (15, "建议", "HistoryPage 无分页性能风险",
     "HistoryPage 一次性加载所有会话，无分页或虚拟滚动，\n对话量大时渲染和搜索性能下降",
     "大量 DOM 节点同时存在会降低渲染性能，\n用户在长列表中滚动可能出现卡顿。",
     "实现虚拟滚动（react-window）或分页加载（每页 20 条），\n搜索时仅渲染匹配结果。"),
    (16, "建议", "流式滚动用 instant 非 smooth",
     "ChatPage 流式期间使用 scrollIntoView({behavior:'instant'})\n而非 'smooth'，消息区域频繁跳动",
     "instant 滚动是瞬时跳变，用户视线跟不上新内容生成速度，\n产生视觉不适。",
     "实现智能节流：每 500ms 最多滚动一次，或使用\nscroll-behavior: smooth 的全局 CSS 作为回退。"),
])

# Remaining slides
s_journey()
s_edge()
s_cognitive()
s_roadmap()
s_appendix()

# Save
os.makedirs(DIR, exist_ok=True)
prs.save(OUT)
print(f"PPT saved to: {OUT}")
print(f"Total slides: {len(prs.slides)}")
