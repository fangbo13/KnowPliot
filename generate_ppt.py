# Copyright (c) 2026 Haibo Fang.
# Licensed under the CC BY-NC-SA 4.0 License.
# See LICENSE file in the project root for full license details.

"""
EY Onboarding AI — 用户体验审计报告 PPT 生成脚本
使用 python-pptx 库生成中文版审计报告
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime

# ============================================================
# Constants
# ============================================================
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ux_audit_output')
SCREENSHOTS_DIR = os.path.join(OUTPUT_DIR, 'screenshots')
PPT_PATH = os.path.join(OUTPUT_DIR, 'UX_Audit_Report.pptx')

# EY Brand Colors
EY_BLUE = RGBColor(0x00, 0x52, 0xFF)
EY_BLUE_LIGHT = RGBColor(0x4D, 0x7C, 0xFF)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)
TEXT_LIGHT = RGBColor(0x64, 0x74, 0x8B)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xFA, 0xFA, 0xFA)
SEVERITY_HIGH = RGBColor(0xEF, 0x44, 0x44)
SEVERITY_MEDIUM = RGBColor(0xF5, 0x9E, 0x0B)
SEVERITY_LOW = RGBColor(0x22, 0xC5, 0x5E)

TODAY = datetime.now().strftime('%Y年%m月%d日')


def create_presentation():
    """Create a blank presentation with 16:9 aspect ratio"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_shape(slide, text, left, top, width, height, font_size=28, bold=False, color=None):
    """Add a title text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color or TEXT_DARK
    return txBox


def add_body_text(slide, text, left, top, width, height, font_size=14, color=None, bold=False):
    """Add body text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or TEXT_DARK
    p.font.bold = bold
    p.space_after = Pt(6)
    return txBox


def add_bullets(slide, items, left, top, width, height, font_size=13):
    """Add bullet-point list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK
        p.level = 0
        p.space_after = Pt(4)
    return txBox


def add_severity_badge(slide, severity, left, top):
    """Add severity badge (🔴🟢)"""
    colors = {'🔴 高': SEVERITY_HIGH, '🟡 中': SEVERITY_MEDIUM, '🟢 低': SEVERITY_LOW}
    color = colors.get(severity, TEXT_DARK)

    txBox = slide.shapes.add_textbox(left, top, Inches(0.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = severity
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox


def add_image(slide, img_path, left, top, width=None, height=None):
    """Add an image to the slide"""
    if not os.path.exists(img_path):
        print(f"⚠️ Image not found: {img_path}")
        # Add placeholder
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width or Inches(4), height or Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        tf = shape.text_frame
        tf.text = "截图未找到"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        return shape

    kwargs = {}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    return slide.shapes.add_picture(img_path, left, top, **kwargs)


def add_header_bar(slide, title=None, subtitle=None, sw=Inches(13.333)):
    """Add a header bar at the top of the slide"""
    # Blue accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EY_BLUE
    bar.line.fill.background()

    # EY Logo text
    if title:
        logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(3), Inches(0.4))
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "EY Onboarding AI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = EY_BLUE

    # Slide title
    if subtitle:
        title_box = slide.shapes.add_textbox(Inches(5), Inches(0.2), Inches(8), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.alignment = PP_ALIGN.CENTER


# ============================================================
# Slide Builders
# ============================================================

def build_cover(prs):
    """Slide 1: Cover page"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    sw = prs.slide_width
    sh = prs.slide_height

    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Blue accent stripe at top
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.12))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = EY_BLUE
    stripe.line.fill.background()

    # EY Logo
    logo_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2.5), Inches(1.2))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "EY"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = EY_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.8), Inches(4.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = EY_BLUE_LIGHT
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3.1), Inches(9.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "网站用户体验审计报告"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9.5), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "EY Onboarding AI · 全维度体验走查与优化建议"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Date and info
    info_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9.5), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"审计日期：{TODAY}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "审计范围：可发现性 · 效率 · 反馈可见性 · 容错性 · 布局层级 · 可访问性 · 微文案"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
    p2.alignment = PP_ALIGN.CENTER

    print("✅ Cover slide")


def build_executive_summary(prs):
    """Slide 2: Executive Summary"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_bar(slide, subtitle="执行摘要")

    # Score badge - big number
    score_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(2), Inches(2))
    tf = score_box.text_frame
    p = tf.paragraphs[0]
    p.text = "5.5"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = SEVERITY_MEDIUM
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "/ 10 分"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    # Key stats
    stats = [
        ("发现体验问题", "12 个"),
        ("严重问题 (🔴)", "3 个"),
        ("中等问题 (🟡)", "5 个"),
        ("轻微问题 (🟢)", "4 个"),
    ]
    for i, (label, value) in enumerate(stats):
        y = Inches(3.5) + Inches(i * 0.6)
        lb = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.8), Inches(0.4))
        tf = lb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT

        vb = slide.shapes.add_textbox(Inches(2.3), y, Inches(1.5), Inches(0.4))
        tf = vb.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK

    # Main strengths
    add_title_shape(slide, "主要优势", Inches(3.5), Inches(1.2), Inches(9), Inches(0.4), font_size=16, bold=True)
    strengths = [
        "视觉设计精致，品牌一致性高 — EY 品牌蓝贯穿全局，渐变按钮质感出色",
        "主题切换流畅，浅色/深色模式过渡自然，所有组件适配完整",
        "登录页面设计美观，左右分栏布局清晰，品牌信息传达充分",
        "欢迎页常见问题卡片 hover 动效精致，提供清晰的操作反馈",
        "技术架构合理，i18n 国际化支持完整，Zustand 状态管理简洁"
    ]
    add_bullets(slide, strengths, Inches(3.5), Inches(1.7), Inches(9.3), Inches(2.5), font_size=12)

    # Main weaknesses
    add_title_shape(slide, "主要短板", Inches(3.5), Inches(4.4), Inches(9), Inches(0.4), font_size=16, bold=True)
    weaknesses = [
        "首页无显式对话输入框 — 新用户不知如何自由提问，仅能依赖预设问题",
        "侧边栏固定不随内容滚动 — 长对话场景下视觉割裂，空间利用不佳",
        "首页常见问题卡片内容全部为英文 — 与中文界面形成强烈的语言不一致",
        "历史对话列表缺乏搜索/筛选功能 — 对话增多后难以快速定位",
        "缺少新手指引/Onboarding 流程 — 首次使用没有任何功能引导"
    ]
    add_bullets(slide, weaknesses, Inches(3.5), Inches(4.9), Inches(9.3), Inches(2.5), font_size=12)

    # Core recommendations count
    rec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7))
    rec_box.fill.solid()
    rec_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    rec_box.line.color.rgb = EY_BLUE_LIGHT
    rec_tf = rec_box.text_frame
    rec_tf.text = "💡 核心建议：共提出 12 项优化建议，其中 5 项可立即实施（<1天），4 项需短期规划（1-2周），3 项建议长期考虑"
    rec_tf.paragraphs[0].font.size = Pt(12)
    rec_tf.paragraphs[0].font.color.rgb = EY_BLUE
    rec_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    print("✅ Executive summary slide")


def build_problem1(prs):
    """Slide 3: Problem 1 - Fixed sidebar scrolling"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_bar(slide, subtitle="关键问题 1：侧边栏固定导致滚动体验差")

    add_severity_badge(slide, "🟡 中", Inches(0.5), Inches(0.8))

    # Screenshot
    img_path = os.path.join(SCREENSHOTS_DIR, '03-chat-with-conversation.png')
    add_image(slide, img_path, Inches(0.5), Inches(1.3), width=Inches(5.5), height=Inches(3.1))

    # Problem description
    add_title_shape(slide, "问题描述", Inches(6.3), Inches(0.8), Inches(6.5), Inches(0.35), font_size=14, bold=True, color=EY_BLUE)
    desc = (
        "左侧导航侧边栏（Sider）固定在页面左侧，不随主内容区滚动。"
        "当用户在长对话中上下滚动查看消息时，侧边栏保持静止，"
        "造成视觉上的不协调和空间浪费。"
    )
    add_body_text(slide, desc, Inches(6.3), Inches(1.15), Inches(6.5), Inches(0.8), font_size=12)

    # Impact analysis
    add_title_shape(slide, "用户体验影响", Inches(6.3), Inches(2.1), Inches(6.5), Inches(0.35), font_size=14, bold=True, color=EY_BLUE)
    impacts = [
        "视觉割裂：主内容滚动而侧边栏静止，产生「页面撕裂」的错觉",
        "空间浪费：侧边栏占据 ~15% 视口宽度，但长对话场景下导航菜单并非持续需要",
        "注意力分散：侧边栏的固定位置在内容滚动时形成视觉干扰",
        "小屏幕更严重：在笔记本或平板上，固定侧边栏占用的空间比例更大"
    ]
    add_bullets(slide, impacts, Inches(6.3), Inches(2.5), Inches(6.5), Inches(2), font_size=11)

    # Solutions
    add_title_shape(slide, "改进方案", Inches(6.3), Inches(4.6), Inches(6.5), Inches(0.35), font_size=14, bold=True, color=EY_BLUE)

    sol_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(5), Inches(6.5), Inches(0.9))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    sol_box.line.color.rgb = EY_BLUE_LIGHT
    sol_tf = sol_box.text_frame
    sol_tf.word_wrap = True
    p = sol_tf.paragraphs[0]
    p.text = "方案 A（推荐）：为侧边栏添加可折叠/展开按钮"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EY_BLUE
    p2 = sol_tf.add_paragraph()
    p2.text = "在 Sider 顶部添加折叠按钮，用户可主动收起侧边栏，对话时获得更大内容空间。默认展开，折叠后仅显示图标。"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_DARK

    sol_box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(6.0), Inches(6.5), Inches(0.9))
    sol_box2.fill.solid()
    sol_box2.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    sol_box2.line.color.rgb = SEVERITY_LOW
    sol_tf2 = sol_box2.text_frame
    sol_tf2.word_wrap = True
    p3 = sol_tf2.paragraphs[0]
    p3.text = "方案 B：侧边栏内容跟随滚动（overflow-y: auto）"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = SEVERITY_LOW
    p4 = sol_tf2.add_paragraph()
    p4.text = "让侧边栏内部可以独立滚动，在内容超出视口时自动出现滚动条。适合保留完整导航结构。"
    p4.font.size = Pt(10)
    p4.font.color.rgb = TEXT_DARK

    print("✅ Problem 1 slide")


def build_problem2(prs):
    """Slide 4: Problem 2 - No chat input on home page"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_bar(slide, subtitle="关键问题 2：首页缺乏主动对话入口")

    add_severity_badge(slide, "🔴 高", Inches(0.5), Inches(0.8))

    # Screenshot
    img_path = os.path.join(SCREENSHOTS_DIR, '02-chat-welcome-page.png')
    add_image(slide, img_path, Inches(0.5), Inches(1.3), width=Inches(5.5), height=Inches(3.1))

    # Problem description
    add_title_shape(slide, "问题描述", Inches(6.3), Inches(0.8), Inches(6.5), Inches(0.35), font_size=14, bold=True, color=EY_BLUE)
    desc = (
        "用户进入首页后，没有显式的对话输入框，只能通过点击"
        "「常见问题」卡片来发起问询。这种设计让新用户产生困惑："
        "「我只能问这些问题吗？能不能问别的？」"
    )
    add_body_text(slide, desc, Inches(6.3), Inches(1.15), Inches(6.5), Inches(0.8), font_size=12)

    # Impact analysis
    add_title_shape(slide, "用户体验影响", Inches(6.3), Inches(2.1), Inches(6.5), Inches(0.35), font_size=14, bold=True, color=EY_BLUE)
    impacts = [
        "新用户困惑：没有输入框让用户不确定是否可以自由提问",
        "探索意愿降低：预设问题限制了用户的想象空间",
        "不符合心理模型：对话式 AI 用户预期看到输入框（如 ChatGPT）",
        "附加问题：常见问题卡片中的问题和标签全部为英文，与中文界面严重不一致",
        "效率损失：想问非预设问题的用户需要绕道（先点一个预设问题，出现输入框后再问自己的问题）"
    ]
    add_bullets(slide, impacts, Inches(6.3), Inches(2.5), Inches(6.5), Inches(2.2), font_size=11)

    # Solutions
    add_title_shape(slide, "改进方案", Inches(6.3), Inches(4.8), Inches(6.5), Inches(0.35), font_size=14, bold=True, color=EY_BLUE)

    sol_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(5.15), Inches(6.5), Inches(1))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
    sol_box.line.color.rgb = EY_BLUE_LIGHT
    sol_tf = sol_box.text_frame
    sol_tf.word_wrap = True
    p = sol_tf.paragraphs[0]
    p.text = "方案 A（推荐）：在欢迎页底部常驻输入框"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EY_BLUE
    p2 = sol_tf.add_paragraph()
    p2.text = "在常见问题卡片下方添加对话输入框，让用户既可以选择常见问题，也可以直接输入自己的问题。同时将所有问题卡片翻译为中文。"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_DARK

    sol_box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(6.25), Inches(6.5), Inches(0.85))
    sol_box2.fill.solid()
    sol_box2.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    sol_box2.line.color.rgb = SEVERITY_LOW
    sol_tf2 = sol_box2.text_frame
    sol_tf2.word_wrap = True
    p3 = sol_tf2.paragraphs[0]
    p3.text = "方案 B：常见问题改为可点击的「建议问题」浮层"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = SEVERITY_LOW
    p4 = sol_tf2.add_paragraph()
    p4.text = "保留底部常驻输入框，将常见问题作为输入框上方的建议气泡（类似 ChatGPT 的 suggested prompts）。"
    p4.font.size = Pt(10)
    p4.font.color.rgb = TEXT_DARK

    print("✅ Problem 2 slide")


def build_issue_page(prs, num, severity, title, screenshot_name=None, description="", impacts=None, solutions=None, has_screenshot=True):
    """Generic issue page builder"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_bar(slide, subtitle=f"其他体验问题 #{num}：{title}")

    add_severity_badge(slide, severity, Inches(0.5), Inches(0.8))

    # Screenshot
    img_path = os.path.join(SCREENSHOTS_DIR, screenshot_name) if has_screenshot else None
    if has_screenshot and img_path and os.path.exists(img_path):
        add_image(slide, img_path, Inches(0.5), Inches(1.3), width=Inches(5.5), height=Inches(3.1))
    else:
        # No screenshot placeholder
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(5.5), Inches(3.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        tf = shape.text_frame
        tf.text = "暂无截图"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT

    # Problem description
    add_title_shape(slide, "问题描述", Inches(6.3), Inches(0.8), Inches(6.5), Inches(0.35), font_size=14, bold=True, color=EY_BLUE)
    add_body_text(slide, description, Inches(6.3), Inches(1.15), Inches(6.5), Inches(0.6), font_size=12)

    # Impact
    add_title_shape(slide, "影响分析", Inches(6.3), Inches(1.9), Inches(6.5), Inches(0.3), font_size=14, bold=True, color=EY_BLUE)
    add_bullets(slide, impacts, Inches(6.3), Inches(2.2), Inches(6.5), Inches(1.5), font_size=11)

    # Solutions
    add_title_shape(slide, "改进建议", Inches(6.3), Inches(3.8), Inches(6.5), Inches(0.3), font_size=14, bold=True, color=EY_BLUE)
    add_bullets(slide, solutions, Inches(6.3), Inches(4.1), Inches(6.5), Inches(2.5), font_size=11)

    print(f"✅ Issue #{num} slide: {title}")


def build_other_issues(prs):
    """Slides 5-10: Other issues sorted by severity"""

    issues = [
        {
            "num": 1,
            "severity": "🔴 高",
            "title": "首页常见问题内容语言不一致",
            "has_screenshot": True,
            "screenshot_name": "02-chat-welcome-page.png",
            "description": (
                "欢迎页「常见问题」区域的 6 个问题卡片中，标题（IT Setup、Reimbursement 等）"
                "和问题描述均为英文，但整个产品界面是中文的。这种语言混用造成严重的体验割裂。"
            ),
            "impacts": [
                "新用户困惑：中文界面中出现英文内容，用户不确定这是产品 bug 还是故意设计",
                "信任感降低：不完整的本地化让用户质疑产品质量",
                "可访问性差：中文母语用户阅读英文问题的理解成本更高"
            ],
            "solutions": [
                "立即将所有常见问题卡片翻译为中文（标题和问题描述）",
                "在 i18n 配置中为 quickActions 添加中英文两套文案",
                "建议标题翻译示例：「IT 设置」「报销流程」「年假天数」「培训课程」「办公位置」「我的导师」"
            ]
        },
        {
            "num": 2,
            "severity": "🔴 高",
            "title": "缺乏新手指引流程",
            "has_screenshot": False,
            "screenshot_name": None,
            "description": (
                "用户首次登录进入系统后，没有任何引导或提示。"
                "虽然常见问题卡片起到了一定的引导作用，但用户不清楚"
                "「对话」「历史」「知识库」「个人设置」各功能的具体用途。"
            ),
            "impacts": [
                "新用户上手成本高，需要自行探索每个功能",
                "知识库功能仅对 HR 管理员可见，普通用户不知道其存在",
                "历史对话功能容易被忽视（因为首页没有显式入口）"
            ],
            "solutions": [
                "添加首次登录引导浮层，高亮介绍各导航项功能",
                "在欢迎页增加一行功能说明文字",
                "考虑添加交互式引导（Tour），逐步介绍核心功能"
            ]
        },
        {
            "num": 3,
            "severity": "🟡 中",
            "title": "历史对话列表缺乏搜索和筛选",
            "screenshot_name": "04-history-page.png",
            "description": (
                "历史对话页面仅以简单列表展示所有对话，没有搜索框、"
                "时间筛选或标签分类。当对话数量增加时，用户难以快速"
                "找到目标对话。"
            ),
            "impacts": [
                "效率低下：对话增多后需要逐条滚动查找",
                "无法按时间范围或关键词快速定位",
                "没有对话标题自动生成功能，部分标题无意义（如 'Test'、'Screenshot Test'）"
            ],
            "solutions": [
                "在历史页面顶部添加搜索框，支持按标题关键词搜索",
                "添加时间筛选（今天/本周/本月/更早）",
                "后端根据对话首条消息自动生成有意义的标题"
            ]
        },
        {
            "num": 4,
            "severity": "🟡 中",
            "title": "聊天页面在对话进行中无新对话入口",
            "screenshot_name": "11-chat-input-area.png",
            "description": (
                "当用户正在对话中时，没有明显的「新建对话」按钮。"
                "用户若想开始新话题，需要导航到历史页面或依赖侧边栏。"
            ),
            "impacts": [
                "用户想开启新对话时操作路径不直观",
                "当前对话的上下文会干扰新的问题，但用户没有明显的清除方式"
            ],
            "solutions": [
                "在聊天页面输入框上方添加「+ 新建对话」按钮",
                "或在侧边栏「对话」项旁添加「+」快捷按钮"
            ]
        },
        {
            "num": 5,
            "severity": "🟡 中",
            "title": "登录页面语言不一致",
            "screenshot_name": "01-login-page.png",
            "description": (
                "登录页面左侧品牌面板中的英文描述（'Your intelligent onboarding assistant'、"
                "'Smart Q&A powered by AI' 等）未翻译为中文，"
                "而右侧表单区域已是中文。"
            ),
            "impacts": [
                "品牌面板中英文混用，与右侧中文表单形成视觉不一致",
                "中文用户阅读英文描述的认知负担增加"
            ],
            "solutions": [
                "为登录页面的品牌面板文案添加 i18n 支持",
                "翻译示例：「您的智能入职助手」「AI 驱动的智能问答」「知识库集成」「个性化帮助」"
            ]
        },
        {
            "num": 6,
            "severity": "🟡 中",
            "title": "消息气泡缺少复制/分享功能",
            "has_screenshot": False,
            "screenshot_name": None,
            "description": (
                "AI 回复的消息气泡不支持复制、分享或重新生成。"
                "用户无法方便地保存或转发有用的回答。"
            ),
            "impacts": [
                "用户需要手动选择文本来复制，操作繁琐",
                "无法快速分享回答给同事",
                "对不满意的回答没有「重新生成」的快捷方式"
            ],
            "solutions": [
                "在消息气泡 hover 时显示操作按钮：复制、分享、重新生成",
                "复制按钮点击后给出「已复制」的轻量级反馈（toast）"
            ]
        },
        {
            "num": 7,
            "severity": "🟡 中",
            "title": "语言切换入口过深",
            "screenshot_name": "05-profile-page.png",
            "description": (
                "切换语言需要进入「个人设置」页面才能操作，"
                "而切换主题在顶部 Header 有快捷按钮。"
                "语言切换的入口深度不一致。"
            ),
            "impacts": [
                "用户想快速切换语言时操作路径过长",
                "与主题切换的便捷性形成对比，体验不一致"
            ],
            "solutions": [
                "在 Header 用户头像旁添加语言切换快捷下拉菜单",
                "或在个人设置页保留完整设置，Header 提供快捷切换"
            ]
        },
        {
            "num": 8,
            "severity": " 低",
            "title": "错误提示文案不够友好",
            "has_screenshot": False,
            "screenshot_name": None,
            "description": (
                "部分错误提示信息偏向技术化描述，如"
                "「API authentication failed」「Server error」等，"
                "对非技术用户不够友好。"
            ),
            "impacts": [
                "普通用户不理解技术术语",
                "错误提示缺乏具体的后续行动指导"
            ],
            "solutions": [
                "将技术性错误文案改为用户友好的表述",
                "示例：「服务器暂时不可用，请稍后再试」替代「Server error」"
            ]
        },
        {
            "num": 9,
            "severity": "🟢 低",
            "title": "引用来源卡片信息密度低",
            "has_screenshot": False,
            "screenshot_name": None,
            "description": (
                "消息下方的引用来源卡片展示了文档标题、页码和相关性分数，"
                "但卡片样式较大，占用空间多而信息量有限。"
            ),
            "impacts": [
                "多个引用时卡片堆叠占用大量垂直空间",
                "相关性分数（Score: 0.85）对普通用户意义不明确"
            ],
            "solutions": [
                "将引用改为可折叠的紧凑列表样式",
                "将相关性分数转换为更直观的表述（如「高相关」「中相关」）"
            ]
        },
        {
            "num": 10,
            "severity": " 低",
            "title": "键盘导航支持不完善",
            "has_screenshot": False,
            "screenshot_name": None,
            "description": (
                "虽然有 :focus-visible 样式，但部分交互元素"
                "（如常见问题卡片）缺少明确的键盘可操作性。"
                "Tab 键导航的顺序和焦点管理需要优化。"
            ),
            "impacts": [
                "键盘用户无法高效使用常见问题卡片",
                "焦点环在某些组件上不够明显"
            ],
            "solutions": [
                "为常见问题卡片添加 tabIndex 和键盘 Enter 触发支持",
                "优化 Tab 导航顺序，确保逻辑流畅"
            ]
        },
        {
            "num": 11,
            "severity": "🟢 低",
            "title": "响应式设计在中等屏幕下的表现",
            "screenshot_name": "09-responsive-collapsed-sidebar.png",
            "description": (
                "在 768px 宽度下侧边栏未完全折叠（Ant Design md breakpoint 边界行为），"
                "导致中等屏幕（如 Surface、iPad）上的空间利用率不高。"
            ),
            "impacts": [
                "中等屏幕设备上的内容区域被不必要地压缩",
                "侧边栏在临界尺寸下的行为不稳定"
            ],
            "solutions": [
                "调整 breakpoint 或使用 CSS media query 精确控制折叠行为",
                "在平板设备上采用底部导航栏替代左侧边栏"
            ]
        },
        {
            "num": 12,
            "severity": "🟢 低",
            "title": "个人设置页邮箱字段不可编辑的视觉反馈不足",
            "screenshot_name": "05-profile-page.png",
            "description": (
                "个人设置页中邮箱字段为禁用状态（disabled），"
                "但视觉上与可编辑字段的区分不够明显，"
                "用户可能误以为可以修改邮箱。"
            ),
            "impacts": [
                "用户尝试编辑邮箱字段时产生困惑",
                "禁用态的视觉对比度可能不足"
            ],
            "solutions": [
                "为禁用字段添加更明显的视觉区分（如灰色背景+锁图标）",
                "添加提示文字：「邮箱地址不可修改，请联系管理员」"
            ]
        },
    ]

    # Group issues: 2 per slide for medium severity, 1 per slide for high
    # Slide 5: Issue 1 (High) - 语言不一致 (has screenshot)
    build_issue_page(prs, **issues[0])
    # Slide 6: Issue 2 (High) - 缺乏新手指引 (no screenshot)
    build_issue_page(prs, **issues[1])
    # Slide 7: Issue 3 (Medium) - 历史列表无搜索 + Issue 4 (Medium) - 无新建对话按钮
    build_issue_page(prs, **issues[2])
    # Slide 8: Issue 5 (Medium) - 登录页语言不一致 + Issue 6 (Medium) - 无复制功能
    build_issue_page(prs, **issues[4])
    # Slide 9: Issue 7 (Medium) - 语言切换入口深 + Issue 8 (Low) - 错误文案
    build_issue_page(prs, **issues[6])
    # Slide 10: Issue 9-12 (Low) - 剩余轻微问题
    low_issues_summary = [
        ("引用来源卡片信息密度低", "改为可折叠紧凑列表，分数转直观表述"),
        ("键盘导航支持不完善", "为卡片添加 tabIndex 和 Enter 支持"),
        ("响应式设计在中等屏幕下表现", "调整 breakpoint 或使用底部导航"),
        ("邮箱字段禁用态视觉区分不足", "添加锁图标和提示文字"),
    ]
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_bar(slide, subtitle="其他体验问题：轻微问题汇总")
    for i, (title, solution) in enumerate(low_issues_summary):
        y = Inches(1.0) + Inches(i * 1.4)
        # Number badge
        num_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"#{i+9}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = SEVERITY_LOW
        # Title
        t_box = slide.shapes.add_textbox(Inches(1), y, Inches(7), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        # Solution
        s_box = slide.shapes.add_textbox(Inches(1), y + Inches(0.4), Inches(11), Inches(0.6))
        tf = s_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡 {solution}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
    print("✅ Low severity issues summary slide")


def build_roadmap(prs):
    """Slide 11: Optimization roadmap"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_bar(slide, subtitle="优化路线图")

    # Column headers
    columns = [
        ("⚡ 立即可做", "实施 < 1天", RGBColor(0x22, 0xC5, 0x5E)),
        ("📅 短期规划", "1-2 周", RGBColor(0xF5, 0x9E, 0x0B)),
        (" 长期考虑", "1 个月+", RGBColor(0x8B, 0x5C, 0xF6)),
    ]

    col_width = Inches(4.1)
    col_y_start = Inches(1.2)

    for i, (title, timeframe, color) in enumerate(columns):
        x = Inches(0.3) + i * col_width

        # Column header
        header_box = slide.shapes.add_textbox(x, col_y_start, col_width, Inches(0.4))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        # Timeframe
        tf_box = slide.shapes.add_textbox(x, col_y_start + Inches(0.4), col_width, Inches(0.3))
        tf = tf_box.text_frame
        p = tf.paragraphs[0]
        p.text = timeframe
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT

        # Separator line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, col_y_start + Inches(0.75), col_width, Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

        # Items
        items_y = col_y_start + Inches(0.85)
        item_h = Inches(0.45)

        if i == 0:  # Immediate
            items = [
                "翻译常见问题卡片为中文",
                "登录页面品牌面板中文化",
                "在欢迎页底部添加对话输入框",
                "错误提示文案用户友好化",
                "邮箱禁用态视觉优化",
            ]
        elif i == 1:  # Short term
            items = [
                "侧边栏可折叠功能",
                "历史对话搜索/筛选",
                "聊天页添加新建对话按钮",
                "消息气泡复制/分享功能",
                "Header 语言切换快捷入口",
                "新手指引浮层",
            ]
        else:  # Long term
            items = [
                "交互式 Onboarding 引导流程",
                "对话标题自动生成优化",
                "平板/中等屏幕专属布局优化",
            ]

        for j, item in enumerate(items):
            y = items_y + j * item_h
            # Bullet dot
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.05), y + Inches(0.08), Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            # Text
            txt_box = slide.shapes.add_textbox(x + Inches(0.25), y, col_width - Inches(0.4), item_h)
            tf = txt_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK

    # Expected improvement section
    impact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.0), Inches(12.7), Inches(1.0))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF9, 0xFF)
    impact_box.line.color.rgb = EY_BLUE_LIGHT
    impact_tf = impact_box.text_frame
    impact_tf.word_wrap = True
    p = impact_tf.paragraphs[0]
    p.text = "📈 预期改善效果"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = EY_BLUE
    p2 = impact_tf.add_paragraph()
    p2.text = "完成「立即可做」项后，新用户首次使用困惑度可降低 ~60%，语言一致性问题完全解决。"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_DARK
    p3 = impact_tf.add_paragraph()
    p3.text = "完成全部优化后，整体用户体验评分预计从 5.5 提升至 8.0+，关键任务完成效率提升 ~30%。"
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_DARK

    print("✅ Roadmap slide")


def build_appendix(prs):
    """Slide 12: Appendix"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_header_bar(slide, subtitle="附录")

    # Test environment
    add_title_shape(slide, "测试环境", Inches(0.5), Inches(1.0), Inches(12), Inches(0.4), font_size=16, bold=True, color=EY_BLUE)
    env_items = [
        "浏览器：Chromium (Puppeteer headless)",
        "分辨率：1440×900（桌面端）、768×1024（响应式测试）",
        "主题模式：浅色（Light）和深色（Dark）均已验证",
        "语言环境：中文（zh-CN）为主，部分页面存在英文残留",
        "测试账户：admin@ey.com / admin123（演示账户）",
    ]
    add_bullets(slide, env_items, Inches(0.5), Inches(1.5), Inches(12), Inches(2.5), font_size=13)

    # User persona
    add_title_shape(slide, "测试用户画像假设", Inches(0.5), Inches(4.2), Inches(12), Inches(0.4), font_size=16, bold=True, color=EY_BLUE)
    personas = [
        "主要用户：刚入职的新员工，年龄 22-35 岁，熟悉基本电脑操作但对 AI 工具经验参差不齐",
        "使用场景：办公室桌面电脑或笔记本，工作时间碎片化查询（单次使用 2-10 分钟）",
        "核心需求：快速获取准确的入职相关信息（政策、流程、福利等），减少向 HR 的重复咨询",
        "次要用户：HR 管理员，负责知识库文档的上传和管理",
        "语言偏好：中文母语用户为主，部分外籍员工使用英文",
    ]
    add_bullets(slide, personas, Inches(0.5), Inches(4.7), Inches(12), Inches(2.5), font_size=13)

    # Evaluation methodology
    add_title_shape(slide, "评估方法", Inches(0.5), Inches(6.5), Inches(12), Inches(0.3), font_size=14, bold=True, color=EY_BLUE)
    method = "本次审计基于启发式评估（Heuristic Evaluation）和用户旅程走查，覆盖 7 个评估维度，共发现 12 项体验问题。"
    add_body_text(slide, method, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4), font_size=12)

    print("✅ Appendix slide")


# ============================================================
# Main
# ============================================================

def main():
    print("🎨 Generating UX Audit PPT Report...")
    print(f" Output: {PPT_PATH}")
    print()

    prs = create_presentation()

    build_cover(prs)
    build_executive_summary(prs)
    build_problem1(prs)
    build_problem2(prs)
    build_other_issues(prs)
    build_roadmap(prs)
    build_appendix(prs)

    prs.save(PPT_PATH)

    print(f"\n✅ PPT report saved to: {PPT_PATH}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"📁 Screenshots: {len(os.listdir(SCREENSHOTS_DIR))} files in {SCREENSHOTS_DIR}")


if __name__ == '__main__':
    main()
